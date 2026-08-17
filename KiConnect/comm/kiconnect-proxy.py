from flask import Flask, request, Response, send_from_directory, abort
import requests
import ipaddress
import socket
import os
import re
import json
import time
import sys
import threading
import contextlib
import webbrowser
import shutil
import signal
import base64
from collections import defaultdict
from functools import wraps
from urllib.parse import urlparse, unquote
try:
    from cryptography.hazmat.primitives.ciphers.aead import AESGCM
except ImportError:
    print('[ERROR] "cryptography" package not installed. Run: pip install cryptography')
    sys.exit(1)

port = 5000

app = Flask(__name__)

# ── Directories ───────────────────────────────────────────────────
STATIC_DIR = os.path.realpath(os.path.dirname(os.path.abspath(__file__)))
DATA_DIR   = os.path.join(STATIC_DIR, 'datas')
os.makedirs(DATA_DIR, exist_ok=True)

# Projects (coding-agent feature) point at real, user-chosen folders anywhere
# on the local filesystem, re-confined to that folder on every call (see
# _project_root_for_id / _safe_rel_path). A folder can never be registered
# if it equals or encloses STATIC_DIR/DATA_DIR.
#
# The project-id -> path mapping lives per-account, encrypted at rest, in
# datas/<accountId>/agent_projects.json. The AES key is derived client-side
# from the account password and held in RAM only during an unlocked session.

# ── Strict Origin / Host Check ────────────────────────────────────
ALLOWED_ORIGINS = {
    f'http://localhost:{port}',
    f'http://127.0.0.1:{port}',
}

# ── Max Size (DoS-Protection) ─────────────────────────────────────
MAX_BODY_SIZE  = 50  * 1024 * 1024   # 50 MB for proxy requests
MAX_STORE_SIZE = 100 * 1024 * 1024   # 100 MB pro Storage-entry
app.config['MAX_CONTENT_LENGTH'] = MAX_STORE_SIZE

# ── Storage lock (thread-safe file I/O) ──────────────────────────
# RLock, not Lock: endpoints wrap "load, check, save" as one nested acquire
# by the same thread - a plain Lock would deadlock there.
_store_lock = threading.RLock()

# ── Input-Validating ─────────────────────────────────────────────
# accountId: Timestamp + random-part, z.B. "1718000000000_ab3f7"
# key: config, providers, profiles, folders, chats, current_chat, ...
_SAFE_ID_RE  = re.compile(r'^[A-Za-z0-9_\-]{1,128}$')
_SAFE_KEY_RE = re.compile(r'^[A-Za-z0-9_\-]{1,64}$')

def _valid_id(v):  return bool(_SAFE_ID_RE.match(v or ''))
def _valid_key(v): return bool(_SAFE_KEY_RE.match(v or ''))

def _account_dir(account_id):
    if not _valid_id(account_id): return None
    path = os.path.realpath(os.path.join(DATA_DIR, account_id))
    if not path.startswith(DATA_DIR + os.sep): return None
    return path

def _key_path(account_id, key):
    adir = _account_dir(account_id)
    if not adir: return None
    if not _valid_key(key): return None
    path = os.path.realpath(os.path.join(adir, key + '.json'))
    if not path.startswith(adir + os.sep): return None
    return path

def _registry_path():
    return os.path.join(DATA_DIR, '_registry.json')

# ── Atomarer Schreibvorgang ───────────────────────────────────────
def _atomic_write(target_path, data_bytes):
    """Write atomically via temp file + os.replace(). Falls back to a direct
    write if os.replace() keeps failing due to file locking (AV scanners,
    Nextcloud, ...). Must be called inside _store_lock."""
    tmp = target_path + '.tmp'
    with open(tmp, 'wb') as f:
        f.write(data_bytes)
    for attempt in range(8):
        try:
            os.replace(tmp, target_path)
            return
        except PermissionError:
            time.sleep(0.15 * (attempt + 1))
    # Last-resort fallback: direct write
    with open(target_path, 'wb') as fw:
        fw.write(data_bytes)
    try:
        os.remove(tmp)
    except Exception:
        pass


# ── /store/ - Account registry ───────────────────────────────────
@app.route('/store/', methods=['GET', 'PUT', 'OPTIONS'])
@app.route('/store',  methods=['GET', 'PUT', 'OPTIONS'])
def store_registry():
    if request.method == 'OPTIONS':
        return Response('', 204, headers={**CORS_HEADERS, **SECURITY_HEADERS})
    rpath = _registry_path()

    if request.method == 'GET':
        with _store_lock:
            if not os.path.isfile(rpath):
                return Response('[]', 200, content_type='application/json')
            with open(rpath, 'rb') as f:
                return Response(f.read(), 200, content_type='application/json')


    # PUT - uses shared _atomic_write helper
    body = request.get_data()
    if len(body) > MAX_STORE_SIZE:
        return Response('{"error":"Body too large."}', 413, content_type='application/json')
    try: json.loads(body)
    except Exception:
        return Response('{"error":"Invalid JSON."}', 400, content_type='application/json')
    with _store_lock:
        _atomic_write(rpath, body)
    return Response('{"ok":true}', 200, content_type='application/json')
    


# ── /store/<accountId> - Keys auflisten / ganzes Konto-Verzeichnis löschen ──
@app.route('/store/<account_id>', methods=['GET', 'DELETE', 'OPTIONS'])
def store_list(account_id):
    if request.method == 'OPTIONS':
        return Response('', 204, headers={**CORS_HEADERS, **SECURITY_HEADERS})
    adir = _account_dir(account_id)
    if not adir:
        return Response('{"error":"Invalid account ID."}', 400, content_type='application/json')

    if request.method == 'DELETE':
        # Removes the account's entire directory (all its .json key files,
        # including agent_projects.json etc.) in one go, so no empty leftover
        # folder remains after an account is deleted. Individual /store/<id>/<key>
        # DELETE calls only ever removed files, never the directory itself.
        with _store_lock:
            if os.path.isdir(adir):
                shutil.rmtree(adir, ignore_errors=True)
        return Response('{"ok":true}', 200, content_type='application/json')

    with _store_lock:
        if not os.path.isdir(adir):
            return Response('[]', 200, content_type='application/json')
        keys = [f[:-5] for f in os.listdir(adir)
                if f.endswith('.json') and _valid_key(f[:-5])]
    return Response(json.dumps(keys), 200, content_type='application/json')


# ── /store/<accountId>/<key> - Lesen / Schreiben / Loeschen ──────
@app.route('/store/<account_id>/<key>', methods=['GET', 'PUT', 'DELETE', 'OPTIONS'])
def store_key(account_id, key):
    if request.method == 'OPTIONS':
        return Response('', 204, headers={**CORS_HEADERS, **SECURITY_HEADERS})
    fpath = _key_path(account_id, key)
    if not fpath:
        return Response('{"error":"Invalid ID or key."}', 400, content_type='application/json')

    if request.method == 'GET':
        with _store_lock:
            if not os.path.isfile(fpath):
                return Response('null', 200, content_type='application/json')
            with open(fpath, 'rb') as f:
                return Response(f.read(), 200, content_type='application/json')

    if request.method == 'PUT':
        #  - uses shared _atomic_write helper
        body = request.get_data()
        if len(body) > MAX_STORE_SIZE:
            return Response('{"error":"Body too large."}', 413, content_type='application/json')
        try: json.loads(body)
        except Exception:
            return Response('{"error":"Invalid JSON."}', 400, content_type='application/json')
        adir = _account_dir(account_id)
        with _store_lock:
            os.makedirs(adir, exist_ok=True)
            try:
                _atomic_write(fpath, body)
            except Exception as e:
                return Response(f'{{"error":"Write failed: {e}"}}', 500, content_type='application/json')
        return Response('{"ok":true}', 200, content_type='application/json')
        

    if request.method == 'DELETE':
        with _store_lock:
            if os.path.isfile(fpath):
                os.remove(fpath)
        return Response('{"ok":true}', 200, content_type='application/json')


# ════════════════════════════════════════════════════════════════
#  Agent-API - file operations for the coding-agent feature.
#  Every path is re-resolved and re-confined to its project's registered
#  root on every call - never trust a previously-checked path.
# ════════════════════════════════════════════════════════════════
MAX_AGENT_FILE_SIZE = 5 * 1024 * 1024          # 5 MB per file (read + write)
MAX_AGENT_TREE_ENTRIES = 4000                  # safety cap for huge folders
AGENT_IGNORE_DIRS = {
    '.git', 'node_modules', '__pycache__', '.venv', 'venv', 'env',
    'dist', 'build', '.idea', '.vscode', '.pytest_cache', '.mypy_cache',
    'target', '.next', '.svelte-kit',
}
# Files the agent may never touch - secrets, credentials, sandbox-escape risks.
AGENT_BLOCKED_SUFFIXES = ('.env', '.key', '.pem', '.pfx', '.p12', '.crt')
AGENT_BLOCKED_NAMES = {'.env', 'id_rsa', 'id_ed25519', '.npmrc', '.pypirc'}
_SAFE_PROJECT_ID_RE = re.compile(r'^p_[A-Za-z0-9]{1,64}$')

# ── Agent Session (per-account, in-RAM only) ──────────────────────
# The AES-256 key for a project registry is derived client-side from the
# account password and handed to the server once at unlock. Kept only in
# this in-memory dict, keyed by a random session token - never persisted.
_agent_sessions = {}          # token(str) -> {'accountId': str, 'key': bytes, 'ts': float}
_AGENT_SESSION_TTL = 24 * 3600
_agent_session_lock = threading.Lock()

def _agent_session_or_401():
    """Look up the caller's agent session from the X-Agent-Session header.
    Returns the session dict, or None if missing/unknown/expired."""
    token = request.headers.get('X-Agent-Session', '')
    with _agent_session_lock:
        sess = _agent_sessions.get(token)
        if not sess:
            return None
        if time.time() - sess['ts'] > _AGENT_SESSION_TTL:
            _agent_sessions.pop(token, None)
            return None
        sess['ts'] = time.time()  # sliding expiry on activity
        return sess

def _agent_decrypt(key_bytes, b64_blob):
    """Decrypt a base64(iv[12] || AES-GCM ciphertext+tag) blob - same format
    encryptStr()/encryptObj() in kiconnect.js produce."""
    raw = base64.b64decode(b64_blob)
    iv, ct = raw[:12], raw[12:]
    plaintext = AESGCM(key_bytes).decrypt(iv, ct, None)
    return json.loads(plaintext.decode('utf-8'))

def _agent_encrypt(key_bytes, obj):
    iv = os.urandom(12)
    ct = AESGCM(key_bytes).encrypt(iv, json.dumps(obj).encode('utf-8'), None)
    return base64.b64encode(iv + ct).decode('ascii')

def _load_agent_registry(sess):
    """Read+decrypt this account's project registry (see _key_path())."""
    fpath = _key_path(sess['accountId'], 'agent_projects')
    if not fpath or not os.path.isfile(fpath):
        return {'projects': []}
    with _store_lock:
        with open(fpath, 'r', encoding='utf-8') as f:
            raw = f.read()
    try:
        blob = json.loads(raw) if raw else None
    except ValueError:
        blob = None
    if not blob:
        return {'projects': []}
    data = _agent_decrypt(sess['key'], blob)
    if not isinstance(data.get('projects'), list):
        data['projects'] = []
    return data

def _save_agent_registry(sess, data):
    fpath = _key_path(sess['accountId'], 'agent_projects')
    if not fpath:
        return
    os.makedirs(os.path.dirname(fpath), exist_ok=True)
    blob = _agent_encrypt(sess['key'], data)
    with _store_lock:
        _atomic_write(fpath, json.dumps(blob).encode('utf-8'))

def _is_blocked_root(target):
    """Refuse drive/filesystem roots and anything that equals, encloses, OR
    is nested inside STATIC_DIR/DATA_DIR (e.g. registering DATA_DIR itself,
    a parent of it, or one of its per-account subfolders as an agent
    project would otherwise expose/allow tampering with every account's
    encrypted storage via the agent file-browser/exec endpoints)."""
    drive_root = os.path.realpath(os.path.splitdrive(target)[0] + os.sep) if os.name == 'nt' else '/'
    if target == os.path.realpath(drive_root):
        return True
    for guarded in (STATIC_DIR, DATA_DIR):
        g = os.path.realpath(guarded)
        if target == g or g.startswith(target + os.sep) or target.startswith(g + os.sep):
            return True
    return False

def _project_root_for_id(pid, sess):
    """Resolve a project id to its confined, still-valid real folder path.
    None if the id is unknown, the folder is gone, or it now encloses the
    app's own files (re-checked every call, never cached). Requires an
    unlocked agent session."""
    if not _SAFE_PROJECT_ID_RE.match(pid or ''):
        return None
    registry = _load_agent_registry(sess)
    entry = next((p for p in registry['projects'] if p.get('id') == pid), None)
    if not entry:
        return None
    target = os.path.realpath(entry.get('path') or '')
    if not os.path.isdir(target) or _is_blocked_root(target):
        return None
    return target

def _safe_rel_path(project_dir, rel_path):
    """Resolve+confine a path inside a project dir. Rejects '..', absolute
    paths, null bytes, blocked filenames, and any symlink escape."""
    if rel_path is None:
        return None
    rel_path = rel_path.replace('\\', '/').strip('/')
    if not rel_path or '\x00' in rel_path or '..' in rel_path.split('/'):
        return None
    basename = os.path.basename(rel_path).lower()
    if basename in AGENT_BLOCKED_NAMES or any(basename.endswith(s) for s in AGENT_BLOCKED_SUFFIXES):
        return None
    full = os.path.realpath(os.path.join(project_dir, rel_path))
    if full != project_dir and not full.startswith(project_dir + os.sep):
        return None
    return full

def _agent_error(msg, status=400):
    return Response(json.dumps({'error': msg}), status, content_type='application/json')

def _agent_ok(payload=None, status=200):
    return Response(json.dumps(payload if payload is not None else {'ok': True}),
                    status, content_type='application/json')

# ── /agent/session/unlock - hand the server a per-account agent key ──
# Called once after a normal password login (see kiconnect.js
# unlockAgentSession()). Key is derived client-side from the password + a
# dedicated salt, independent of the config/providers/chats key - a leak of
# one doesn't expose the other.
@app.route('/agent/session/unlock', methods=['POST', 'OPTIONS'])
def agent_session_unlock():
    if request.method == 'OPTIONS':
        return Response('', 204, headers={**CORS_HEADERS, **SECURITY_HEADERS})
    try: body = request.get_json(force=True, silent=True) or {}
    except Exception: body = {}
    account_id = body.get('accountId')
    key_b64 = body.get('key')
    if not _valid_id(account_id) or not key_b64:
        return _agent_error('Invalid request.', 400)
    try:
        key_bytes = base64.b64decode(key_b64)
        if len(key_bytes) != 32:
            raise ValueError('bad key length')
    except Exception:
        return _agent_error('Invalid key.', 400)
    sess = {'accountId': account_id, 'key': key_bytes, 'ts': time.time()}
    try:
        data = _load_agent_registry(sess)
    except Exception:
        # Wrong key for existing ciphertext (e.g. password/account mismatch)
        return _agent_error('Could not unlock project registry.', 400)
    token = os.urandom(24).hex()
    with _agent_session_lock:
        _agent_sessions[token] = sess
    return _agent_ok({'token': token, 'projects': data['projects']})

# ── /agent/session/rekey - password change: re-encrypt under new key ─
# Requires a still-valid OLD session. Re-encrypts the server's own
# already-decrypted registry, never a client-supplied one, so a password
# change can't smuggle in unvalidated project entries.
@app.route('/agent/session/rekey', methods=['POST', 'OPTIONS'])
def agent_session_rekey():
    if request.method == 'OPTIONS':
        return Response('', 204, headers={**CORS_HEADERS, **SECURITY_HEADERS})
    old_sess = _agent_session_or_401()
    if not old_sess:
        return _agent_error('Session expired - please log in again.', 401)
    try: body = request.get_json(force=True, silent=True) or {}
    except Exception: body = {}
    new_key_b64 = body.get('newKey')
    if not new_key_b64:
        return _agent_error('Invalid request.', 400)
    try:
        new_key_bytes = base64.b64decode(new_key_b64)
        if len(new_key_bytes) != 32:
            raise ValueError('bad key length')
    except Exception:
        return _agent_error('Invalid key.', 400)
    with _store_lock:
        registry = _load_agent_registry(old_sess)
        new_sess = {'accountId': old_sess['accountId'], 'key': new_key_bytes, 'ts': time.time()}
        _save_agent_registry(new_sess, registry)
    old_token = request.headers.get('X-Agent-Session', '')
    new_token = os.urandom(24).hex()
    with _agent_session_lock:
        _agent_sessions.pop(old_token, None)
        _agent_sessions[new_token] = new_sess
    return _agent_ok({'token': new_token, 'projects': registry['projects']})

# ── /agent/session/lock - explicit logout for the agent session ──────
@app.route('/agent/session/lock', methods=['POST', 'OPTIONS'])
def agent_session_lock():
    if request.method == 'OPTIONS':
        return Response('', 204, headers={**CORS_HEADERS, **SECURITY_HEADERS})
    token = request.headers.get('X-Agent-Session', '')
    with _agent_session_lock:
        _agent_sessions.pop(token, None)
    return _agent_ok()

# ── /agent/browse - list real OS folders, for the project folder picker ──
def _list_subdirs(path):
    names = []
    try:
        with os.scandir(path) as it:
            for e in it:
                try:
                    if e.is_dir(follow_symlinks=False) and not e.name.startswith('.'):
                        names.append(e.name)
                except OSError:
                    continue
    except OSError:
        return None
    return sorted(names, key=str.lower)

def _list_files(path):
    """Lists files directly inside `path` whose extension the KB indexer
    can handle (KB_SUPPORTED_EXT, defined further below in the KB section -
    looked up at call time, so this forward reference is fine: by the time
    a request actually runs, the whole module has finished loading). Used
    by the KB "add individual files" picker (files=1 on /agent/browse)."""
    out = []
    try:
        with os.scandir(path) as it:
            for e in it:
                try:
                    if not e.is_file(follow_symlinks=False) or e.name.startswith('.'):
                        continue
                    ext = os.path.splitext(e.name)[1].lower()
                    if ext not in KB_SUPPORTED_EXT:
                        continue
                    try:
                        if e.stat().st_size > MAX_KB_FILE_SIZE:
                            continue
                    except OSError:
                        continue
                    out.append(e.name)
                except OSError:
                    continue
    except OSError:
        return None
    return sorted(out, key=str.lower)

@app.route('/agent/browse', methods=['GET', 'OPTIONS'])
def agent_browse():
    if request.method == 'OPTIONS':
        return Response('', 204, headers={**CORS_HEADERS, **SECURITY_HEADERS})
    raw = (request.args.get('path') or '').strip()
    want_files = (request.args.get('files') or '') == '1'
    target = os.path.realpath(raw) if raw else os.path.realpath(os.path.expanduser('~'))
    if not os.path.isdir(target):
        return _agent_error('Folder not found or not accessible.', 404)
    names = _list_subdirs(target)
    if names is None:
        return _agent_error('No access to this folder.', 403)
    parent = os.path.dirname(target)
    if parent == target or not os.path.isdir(parent):
        parent = None
    shortcuts = [{'label': 'Home', 'path': os.path.realpath(os.path.expanduser('~'))}]
    if os.name == 'nt':
        import string
        for letter in string.ascii_uppercase:
            drive = f'{letter}:\\'
            if os.path.exists(drive):
                shortcuts.append({'label': drive, 'path': os.path.realpath(drive)})
    else:
        shortcuts.append({'label': '/', 'path': '/'})
    out = {
        'path': target,
        'parent': parent,
        'entries': [{'name': n, 'path': os.path.join(target, n)} for n in names],
        'shortcuts': shortcuts,
    }
    if want_files:
        file_names = _list_files(target)
        out['files'] = [{'name': n, 'path': os.path.join(target, n)} for n in (file_names or [])]
    return _agent_ok(out)

# ── /agent/projects - list / register project folders ────────────
@app.route('/agent/projects', methods=['GET', 'POST', 'OPTIONS'])
def agent_projects():
    if request.method == 'OPTIONS':
        return Response('', 204, headers={**CORS_HEADERS, **SECURITY_HEADERS})
    sess = _agent_session_or_401()
    if not sess:
        return _agent_error('Session expired - please log in again.', 401)

    if request.method == 'GET':
        registry = _load_agent_registry(sess)
        projects = []
        for p in registry['projects']:
            target = os.path.realpath(p.get('path') or '')
            projects.append({
                'id': p.get('id'), 'name': p.get('name'), 'path': target,
                'missing': not (os.path.isdir(target) and not _is_blocked_root(target)),
                'shell': bool(p.get('shell')),
            })
        return _agent_ok({'projects': projects})

    # POST - register an existing (or newly created) real folder as a project
    try: body = request.get_json(force=True, silent=True) or {}
    except Exception: body = {}
    name = (body.get('name') or '').strip()
    raw_path = (body.get('path') or '').strip()
    create = bool(body.get('create'))
    if not name or len(name) > 64:
        return _agent_error('Invalid project name (1-64 characters).')
    if not raw_path:
        return _agent_error('Please provide a folder path.')
    target = os.path.realpath(raw_path)
    if create and not os.path.isdir(target):
        try:
            os.makedirs(target, exist_ok=True)
        except OSError as e:
            return _agent_error(f'Could not create folder: {e}')
    if not os.path.isdir(target):
        return _agent_error('Folder not found.', 404)
    if _is_blocked_root(target):
        return _agent_error('This folder is not allowed for security reasons (drive/system root or the app\'s own folder).', 403)
    with _store_lock:
        registry = _load_agent_registry(sess)
        if any(os.path.realpath(p.get('path') or '') == target for p in registry['projects']):
            return _agent_error('This folder is already registered as a project.', 409)
        pid = 'p_' + os.urandom(10).hex()
        # Shell execution is OFF by default for every new project - it's an
        # explicit, separate opt-in per project (see /agent/projects/<id>/shell).
        registry['projects'].append({'id': pid, 'name': name, 'path': target, 'shell': False})
        _save_agent_registry(sess, registry)
    return _agent_ok({'id': pid, 'name': name, 'path': target})

# ── /agent/projects/<id>/shell - enable/disable shell command execution ──
# Kept as its own explicit endpoint (rather than folded into general project
# settings) so turning this on is always a distinct, deliberate action, and
# so /agent/exec below can independently double-check it server-side rather
# than trusting whatever the frontend currently displays.
@app.route('/agent/projects/<pid>/shell', methods=['PUT', 'OPTIONS'])
def agent_set_shell(pid):
    if request.method == 'OPTIONS':
        return Response('', 204, headers={**CORS_HEADERS, **SECURITY_HEADERS})
    sess = _agent_session_or_401()
    if not sess:
        return _agent_error('Session expired - please log in again.', 401)
    try: body = request.get_json(force=True, silent=True) or {}
    except Exception: body = {}
    enabled = bool(body.get('enabled'))
    with _store_lock:
        registry = _load_agent_registry(sess)
        entry = next((p for p in registry['projects'] if p.get('id') == pid), None)
        if not entry:
            return _agent_error('Project not found.', 404)
        entry['shell'] = enabled
        _save_agent_registry(sess, registry)
    return _agent_ok({'id': pid, 'shell': enabled})

# ── /agent/projects/<id>/path - change a project's target folder ─────
# Re-points an already-registered project without deleting/recreating it.
# Runs through the same validation as creating a new project (existence,
# 'create' opt-in, blocked-root check, no duplicate-path collision).
@app.route('/agent/projects/<pid>/path', methods=['PUT', 'OPTIONS'])
def agent_set_project_path(pid):
    if request.method == 'OPTIONS':
        return Response('', 204, headers={**CORS_HEADERS, **SECURITY_HEADERS})
    sess = _agent_session_or_401()
    if not sess:
        return _agent_error('Session expired - please log in again.', 401)
    try: body = request.get_json(force=True, silent=True) or {}
    except Exception: body = {}
    raw_path = (body.get('path') or '').strip()
    create = bool(body.get('create'))
    if not raw_path:
        return _agent_error('Please provide a folder path.')
    target = os.path.realpath(raw_path)
    if create and not os.path.isdir(target):
        try:
            os.makedirs(target, exist_ok=True)
        except OSError as e:
            return _agent_error(f'Could not create folder: {e}')
    if not os.path.isdir(target):
        return _agent_error('Folder not found.', 404)
    if _is_blocked_root(target):
        return _agent_error('This folder is not allowed for security reasons (drive/system root or the app\'s own folder).', 403)
    with _store_lock:
        registry = _load_agent_registry(sess)
        entry = next((p for p in registry['projects'] if p.get('id') == pid), None)
        if not entry:
            return _agent_error('Project not found.', 404)
        if any(os.path.realpath(p.get('path') or '') == target for p in registry['projects'] if p.get('id') != pid):
            return _agent_error('This folder is already registered as another project.', 409)
        entry['path'] = target
        _save_agent_registry(sess, registry)
    return _agent_ok({'id': pid, 'path': target})

# ── /agent/exec/<id> - run a shell command inside the project folder ──
# Off by default, gated behind the per-project 'shell' flag above.
#
# NOT a hard security boundary (no container/VM) - the command runs as the
# same OS user as the proxy, confined to the project folder only by
# convention (cwd); `cd ..` or an absolute path can still escape. Best-effort
# hardening applied instead:
#   1. Minimal, secret-free environment (doesn't inherit the proxy's env).
#   2. POSIX resource limits (CPU, memory, proc count, open files, no core
#      dumps) via preexec_fn, guarding against fork bombs/runaway loops.
#      Linux/macOS only.
#   3. Own process group (os.setsid) so a timeout can kill the whole subtree.
#   4. Best-effort network isolation via `unshare --net` (Linux only, needs
#      unprivileged user namespaces; probed once, reported back via
#      `networkIsolated` rather than assumed on).
MAX_EXEC_OUTPUT = 200_000  # chars, per stream
EXEC_TIMEOUT_SECONDS = 45
SANDBOX_CPU_SECONDS = EXEC_TIMEOUT_SECONDS + 5   # hard CPU ceiling, just above the wall-clock timeout
SANDBOX_MEM_BYTES = 1024 * 1024 * 1024           # 1 GB address space per command
SANDBOX_MAX_PROCS = 64                            # caps fork bombs
SANDBOX_MAX_OPEN_FILES = 256

_unshare_net_available = None  # cached probe result, see _probe_unshare_net()

def _probe_unshare_net():
    """Checks once whether `unshare --user --map-root-user --net` works
    here. Cached for the process lifetime; never raises."""
    global _unshare_net_available
    if _unshare_net_available is not None:
        return _unshare_net_available
    if os.name == 'nt' or not shutil.which('unshare'):
        _unshare_net_available = False
        return False
    try:
        import subprocess as _sp
        probe = _sp.run(
            ['unshare', '--user', '--map-root-user', '--net', '--', 'true'],
            capture_output=True, timeout=5,
        )
        _unshare_net_available = (probe.returncode == 0)
    except Exception:
        _unshare_net_available = False
    return _unshare_net_available

def _sandbox_env():
    """Minimal environment for sandboxed commands - enough for common
    toolchains (node, python, npm, git) without inheriting the proxy's
    own env vars/secrets."""
    keep = ('PATH', 'HOME', 'LANG', 'LC_ALL', 'TERM', 'TMPDIR', 'USER',
            'SHELL', 'NVM_DIR', 'PYTHONIOENCODING', 'SYSTEMROOT', 'APPDATA')
    env = {k: os.environ[k] for k in keep if k in os.environ}
    env.setdefault('PATH', os.environ.get('PATH', '/usr/local/bin:/usr/bin:/bin'))
    return env

def _sandbox_preexec():
    """Runs in the child right after fork(), before exec - POSIX only.
    Detaches into a new process group and applies resource limits
    best-effort (a rejected limit is skipped, not fatal)."""
    try:
        os.setsid()
    except OSError:
        pass
    import resource
    for res, limits in (
        (resource.RLIMIT_CPU,   (SANDBOX_CPU_SECONDS, SANDBOX_CPU_SECONDS)),
        (resource.RLIMIT_AS,    (SANDBOX_MEM_BYTES, SANDBOX_MEM_BYTES)),
        (resource.RLIMIT_NPROC, (SANDBOX_MAX_PROCS, SANDBOX_MAX_PROCS)),
        (resource.RLIMIT_NOFILE,(SANDBOX_MAX_OPEN_FILES, SANDBOX_MAX_OPEN_FILES)),
        (resource.RLIMIT_CORE,  (0, 0)),
    ):
        try:
            resource.setrlimit(res, limits)
        except (ValueError, OSError):
            pass

@app.route('/agent/exec/<pid>', methods=['POST', 'OPTIONS'])
def agent_exec(pid):
    if request.method == 'OPTIONS':
        return Response('', 204, headers={**CORS_HEADERS, **SECURITY_HEADERS})
    sess = _agent_session_or_401()
    if not sess:
        return _agent_error('Session expired - please log in again.', 401)
    registry = _load_agent_registry(sess)
    entry = next((p for p in registry['projects'] if p.get('id') == pid), None)
    if not entry:
        return _agent_error('Project not found.', 404)
    if not entry.get('shell'):
        return _agent_error('Shell commands are not enabled for this project (⚙ Agent Settings).', 403)
    pdir = _project_root_for_id(pid, sess)
    if not pdir:
        return _agent_error('Project folder not found - was it moved or deleted?', 404)

    try: body = request.get_json(force=True, silent=True) or {}
    except Exception: body = {}
    command = (body.get('command') or '').strip()
    if not command:
        return _agent_error('No command provided.')
    cwd = pdir
    if body.get('cwd'):
        resolved = _safe_rel_path(pdir, body['cwd'])
        if not resolved or not os.path.isdir(resolved):
            return _agent_error('Invalid working directory.')
        cwd = resolved

    import subprocess
    use_posix_limits = (os.name != 'nt')
    net_isolated = use_posix_limits and _probe_unshare_net()
    wrapped_command = ('unshare --user --map-root-user --net -- ' + command) if net_isolated else command

    popen_kwargs = dict(
        cwd=cwd, shell=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE,
        text=True, errors='replace', env=_sandbox_env(),
    )
    if use_posix_limits:
        popen_kwargs['preexec_fn'] = _sandbox_preexec
    else:
        popen_kwargs['creationflags'] = subprocess.CREATE_NEW_PROCESS_GROUP

    try:
        proc = subprocess.Popen(wrapped_command, **popen_kwargs)
    except OSError as e:
        return _agent_error(f'Could not start command: {e}')

    timed_out = False
    try:
        stdout, stderr = proc.communicate(timeout=EXEC_TIMEOUT_SECONDS)
        exit_code = proc.returncode
    except subprocess.TimeoutExpired:
        timed_out, exit_code = True, None
        try:
            if use_posix_limits:
                os.killpg(proc.pid, signal.SIGKILL)  # kill the whole group we detached into via setsid()
            else:
                proc.kill()
        except OSError:
            pass
        stdout, stderr = proc.communicate()

    truncated = len(stdout) > MAX_EXEC_OUTPUT or len(stderr) > MAX_EXEC_OUTPUT
    return _agent_ok({
        'command': command, 'exitCode': exit_code, 'timedOut': timed_out, 'truncated': truncated,
        'stdout': stdout[:MAX_EXEC_OUTPUT], 'stderr': stderr[:MAX_EXEC_OUTPUT],
        'sandboxed': use_posix_limits, 'networkIsolated': net_isolated,
    })



# Only removes the app's link to the folder - the folder itself is never
# deleted, it's real user-owned data outside the app.
@app.route('/agent/projects/<pid>', methods=['DELETE', 'OPTIONS'])
def agent_delete_project(pid):
    if request.method == 'OPTIONS':
        return Response('', 204, headers={**CORS_HEADERS, **SECURITY_HEADERS})
    sess = _agent_session_or_401()
    if not sess:
        return _agent_error('Session expired - please log in again.', 401)
    with _store_lock:
        registry = _load_agent_registry(sess)
        remaining = [p for p in registry['projects'] if p.get('id') != pid]
        if len(remaining) == len(registry['projects']):
            return _agent_error('Project not found.', 404)
        registry['projects'] = remaining
        _save_agent_registry(sess, registry)
    return _agent_ok()

# ── /agent/search/<id> - grep-style text search across a project ──
@app.route('/agent/search/<pid>', methods=['GET', 'OPTIONS'])
def agent_search(pid):
    if request.method == 'OPTIONS':
        return Response('', 204, headers={**CORS_HEADERS, **SECURITY_HEADERS})
    sess = _agent_session_or_401()
    if not sess:
        return _agent_error('Session expired - please log in again.', 401)
    pdir = _project_root_for_id(pid, sess)
    if not pdir:
        return _agent_error('Project folder not found - was it moved or deleted?', 404)
    query = (request.args.get('q') or '').strip()
    if not query:
        return _agent_error('Please provide a search term (q=...).')
    use_regex = request.args.get('regex') == '1'
    case_sensitive = request.args.get('case') == '1'
    scope = (request.args.get('path') or '').strip()
    scope_dir = pdir
    single_file = None  # set below if `path` turns out to name a FILE, not a folder
    note = None
    if scope:
        resolved = _safe_rel_path(pdir, scope)
        if not resolved:
            # Most common cause: model prefixes the project's own folder
            # name onto an already-relative path (paths are always relative
            # to the project root, never including the root's own name).
            return _agent_error(f'Invalid search path "{scope}" - must be a relative path inside the project (no ".." or a leading "/", and not duplicating the project folder\'s own name).')
        if os.path.isfile(resolved):
            # `path` is documented as folder-only, but weaker models
            # sometimes pass a filename instead of calling read_file.
            # Degrade gracefully and just search that one file.
            single_file = resolved
            note = f'"{scope}" is a file, not a folder - `path` normally scopes the search to a subfolder. Searched just that one file. Use read_file to read a specific file directly, or omit `path` to search the whole project.'
        elif not os.path.isdir(resolved):
            return _agent_error(f'Invalid search path "{scope}" - no such file or folder in this project.')
        else:
            scope_dir = resolved
    try:
        pattern = re.compile(query if use_regex else re.escape(query), 0 if case_sensitive else re.IGNORECASE)
    except re.error as e:
        return _agent_error(f'Invalid regular expression: {e}')

    MAX_MATCHES, MAX_FILES = 200, 3000
    matches, files_scanned = [], 0
    file_iter = [(os.path.dirname(single_file), [], [os.path.basename(single_file)])] if single_file else os.walk(scope_dir)
    for root, dirs, files in file_iter:
        if not single_file:
            dirs[:] = sorted(d for d in dirs if d not in AGENT_IGNORE_DIRS and not d.startswith('.'))
        if len(matches) >= MAX_MATCHES or files_scanned >= MAX_FILES:
            break
        for fname in sorted(files):
            if len(matches) >= MAX_MATCHES or files_scanned >= MAX_FILES:
                break
            fpath = os.path.join(root, fname)
            try:
                if os.path.getsize(fpath) > MAX_AGENT_FILE_SIZE:
                    continue
            except OSError:
                continue
            files_scanned += 1
            try:
                with open(fpath, 'r', encoding='utf-8') as f:
                    for lineno, line in enumerate(f, 1):
                        if pattern.search(line):
                            rel = os.path.relpath(fpath, pdir).replace('\\', '/')
                            matches.append({'path': rel, 'line': lineno, 'text': line.strip()[:300]})
                            if len(matches) >= MAX_MATCHES:
                                break
            except (UnicodeDecodeError, OSError):
                continue
    result = {
        'query': query, 'matches': matches,
        'filesScanned': files_scanned, 'truncated': len(matches) >= MAX_MATCHES,
    }
    if note:
        result['note'] = note
    return _agent_ok(result)

# ── /agent/tree/<id> - recursive file listing ─────────────────────
@app.route('/agent/tree/<pid>', methods=['GET', 'OPTIONS'])
def agent_tree(pid):
    if request.method == 'OPTIONS':
        return Response('', 204, headers={**CORS_HEADERS, **SECURITY_HEADERS})
    sess = _agent_session_or_401()
    if not sess:
        return _agent_error('Session expired - please log in again.', 401)
    pdir = _project_root_for_id(pid, sess)
    if not pdir:
        return _agent_error('Project folder not found - was it moved or deleted?', 404)
    entries = []
    count = 0
    for root, dirs, files in os.walk(pdir):
        dirs[:] = sorted(d for d in dirs if d not in AGENT_IGNORE_DIRS and not d.startswith('.'))
        rel_root = os.path.relpath(root, pdir)
        for fname in sorted(files):
            if count >= MAX_AGENT_TREE_ENTRIES:
                break
            rel = fname if rel_root == '.' else f'{rel_root}/{fname}'
            try:
                size = os.path.getsize(os.path.join(root, fname))
            except OSError:
                size = 0
            entries.append({'path': rel.replace('\\', '/'), 'size': size})
            count += 1
        if count >= MAX_AGENT_TREE_ENTRIES:
            break
    return _agent_ok({'project': pid, 'files': entries, 'truncated': count >= MAX_AGENT_TREE_ENTRIES})

# ── /agent/file/<id>/<path> - read / write / delete a file ────────
@app.route('/agent/file/<pid>/<path:rel_path>', methods=['GET', 'PUT', 'DELETE', 'OPTIONS'])
def agent_file(pid, rel_path):
    if request.method == 'OPTIONS':
        return Response('', 204, headers={**CORS_HEADERS, **SECURITY_HEADERS})
    sess = _agent_session_or_401()
    if not sess:
        return _agent_error('Session expired - please log in again.', 401)
    pdir = _project_root_for_id(pid, sess)
    if not pdir:
        return _agent_error('Project folder not found - was it moved or deleted?', 404)
    fpath = _safe_rel_path(pdir, rel_path)
    if not fpath:
        return _agent_error('Invalid file path.')

    if request.method == 'GET':
        if not os.path.isfile(fpath):
            return _agent_error('File not found.', 404)
        if os.path.getsize(fpath) > MAX_AGENT_FILE_SIZE:
            return _agent_error('File too large (>5 MB).', 413)
        with open(fpath, 'rb') as f:
            raw = f.read()
        try:
            return _agent_ok({'path': rel_path, 'content': raw.decode('utf-8'), 'binary': False})
        except UnicodeDecodeError:
            # Not UTF-8 text (PDF, image, zip, ...). Still return the raw
            # bytes as base64 — the frontend can then do format-specific
            # extraction (e.g. pdf.js text extraction for .pdf) instead of
            # just reporting "binary file, can't read it". `content` stays
            # None for backward compatibility with anything that only
            # checks `binary` and expects a plain-text `content` field.
            return _agent_ok({'path': rel_path, 'content': None, 'binary': True, 'content_b64': base64.b64encode(raw).decode('ascii')})

    if request.method == 'PUT':
        try: body = request.get_json(force=True, silent=True) or {}
        except Exception: body = {}
        content = body.get('content', '')
        create_only = bool(body.get('createOnly'))
        if not isinstance(content, str):
            return _agent_error('content must be a string.')
        if len(content.encode('utf-8')) > MAX_AGENT_FILE_SIZE:
            return _agent_error('File too large (>5 MB).', 413)
        if create_only and os.path.exists(fpath):
            return _agent_error('File already exists.', 409)
        with _store_lock:
            os.makedirs(os.path.dirname(fpath), exist_ok=True)
            _atomic_write(fpath, content.encode('utf-8'))
        return _agent_ok({'path': rel_path, 'bytes': len(content.encode('utf-8'))})

    if request.method == 'DELETE':
        with _store_lock:
            if os.path.isfile(fpath):
                os.remove(fpath)
            elif not os.path.exists(fpath):
                return _agent_error('File not found.', 404)
        return _agent_ok({'path': rel_path})

# ── /agent/dir/<id>/<path> - create / delete a directory ──────────
@app.route('/agent/dir/<pid>/<path:rel_path>', methods=['POST', 'DELETE', 'OPTIONS'])
def agent_dir(pid, rel_path):
    if request.method == 'OPTIONS':
        return Response('', 204, headers={**CORS_HEADERS, **SECURITY_HEADERS})
    sess = _agent_session_or_401()
    if not sess:
        return _agent_error('Session expired - please log in again.', 401)
    pdir = _project_root_for_id(pid, sess)
    if not pdir:
        return _agent_error('Project folder not found - was it moved or deleted?', 404)
    dpath = _safe_rel_path(pdir, rel_path)
    if not dpath:
        return _agent_error('Invalid folder path.')

    if request.method == 'POST':
        with _store_lock:
            os.makedirs(dpath, exist_ok=True)
        return _agent_ok({'path': rel_path})

    if request.method == 'DELETE':
        import shutil
        if not os.path.isdir(dpath):
            return _agent_error('Folder not found.', 404)
        with _store_lock:
            shutil.rmtree(dpath, ignore_errors=True)
        return _agent_ok({'path': rel_path})


# ── /agent/move/<id> - move/rename a file or folder ────────────────
# shutil.move (not os.rename) so it also works across directories/filesystems.
@app.route('/agent/move/<pid>', methods=['POST', 'OPTIONS'])
def agent_move(pid):
    if request.method == 'OPTIONS':
        return Response('', 204, headers={**CORS_HEADERS, **SECURITY_HEADERS})
    sess = _agent_session_or_401()
    if not sess:
        return _agent_error('Session expired - please log in again.', 401)
    pdir = _project_root_for_id(pid, sess)
    if not pdir:
        return _agent_error('Project folder not found - was it moved or deleted?', 404)
    try: body = request.get_json(force=True, silent=True) or {}
    except Exception: body = {}
    rel_from = body.get('from')
    rel_to = body.get('to')
    src = _safe_rel_path(pdir, rel_from) if rel_from else None
    dst = _safe_rel_path(pdir, rel_to) if rel_to else None
    if not src or not dst:
        return _agent_error('Invalid from/to path.')
    if not os.path.exists(src):
        return _agent_error('Source not found.', 404)
    if os.path.exists(dst):
        if not bool(body.get('overwrite')):
            return _agent_error('Destination already exists (set overwrite to replace it).', 409)
        with _store_lock:
            if os.path.isdir(dst) and not os.path.islink(dst):
                shutil.rmtree(dst, ignore_errors=True)
            else:
                os.remove(dst)
    with _store_lock:
        os.makedirs(os.path.dirname(dst), exist_ok=True)
        shutil.move(src, dst)
    return _agent_ok({'from': rel_from, 'to': rel_to})


# ── /agent/copy/<id> - copy a file or folder, leaving the original ────
# Mirrors agent_move above, but copies instead of renaming, leaving the
# source untouched. Folder copies are size-capped so a single "copy" call
# can't balloon disk usage with an unbounded recursive copy.
MAX_AGENT_COPY_DIR_BYTES = 200 * 1024 * 1024   # 200 MB cap for a folder copy

def _dir_size(path):
    total = 0
    for root, _dirs, files in os.walk(path):
        for name in files:
            fp = os.path.join(root, name)
            try:
                if not os.path.islink(fp):
                    total += os.path.getsize(fp)
            except OSError:
                pass
    return total

@app.route('/agent/copy/<pid>', methods=['POST', 'OPTIONS'])
def agent_copy(pid):
    if request.method == 'OPTIONS':
        return Response('', 204, headers={**CORS_HEADERS, **SECURITY_HEADERS})
    sess = _agent_session_or_401()
    if not sess:
        return _agent_error('Session expired - please log in again.', 401)
    pdir = _project_root_for_id(pid, sess)
    if not pdir:
        return _agent_error('Project folder not found - was it moved or deleted?', 404)
    try: body = request.get_json(force=True, silent=True) or {}
    except Exception: body = {}
    rel_from = body.get('from')
    rel_to = body.get('to')
    src = _safe_rel_path(pdir, rel_from) if rel_from else None
    dst = _safe_rel_path(pdir, rel_to) if rel_to else None
    if not src or not dst:
        return _agent_error('Invalid from/to path.')
    if not os.path.exists(src):
        return _agent_error('Source not found.', 404)
    if src == dst:
        return _agent_error('Source and destination are the same path.')
    # A folder can't be copied into its own subtree (mirrors the "..\ nested
    # folder" footgun agent.warnSelfNested already warns about client-side).
    if os.path.isdir(src) and (dst == src or dst.startswith(src + os.sep)):
        return _agent_error('Cannot copy a folder into itself or a subfolder of itself.')
    if os.path.isdir(src):
        size = _dir_size(src)
        if size > MAX_AGENT_COPY_DIR_BYTES:
            return _agent_error(f'Folder too large to copy ({size // (1024*1024)} MB > {MAX_AGENT_COPY_DIR_BYTES // (1024*1024)} MB limit).', 413)
    else:
        if os.path.getsize(src) > MAX_AGENT_FILE_SIZE:
            return _agent_error('File too large (>5 MB).', 413)
    if os.path.exists(dst):
        if not bool(body.get('overwrite')):
            return _agent_error('Destination already exists (set overwrite to replace it).', 409)
        with _store_lock:
            if os.path.isdir(dst) and not os.path.islink(dst):
                shutil.rmtree(dst, ignore_errors=True)
            else:
                os.remove(dst)
    with _store_lock:
        os.makedirs(os.path.dirname(dst), exist_ok=True)
        if os.path.isdir(src):
            shutil.copytree(src, dst, symlinks=False)
        else:
            shutil.copy2(src, dst)
    return _agent_ok({'from': rel_from, 'to': rel_to})

# ════════════════════════════════════════════════════════════════
#  Knowledge-Base API ("Wissensbasis" / RAG) - reuses the Agent
#  session above (_agent_session_or_401 / _agent_encrypt / _agent_decrypt
#  / _atomic_write) instead of opening a third independent one, per
#  kiconnect-rag-spec.md section 2: unlocking the project feature also
#  unlocks this one, no extra password round trip for the user.
#
#  Registry: datas/<accountId>/kb_registry.json, encrypted at rest with
#  the SAME key as agent_projects.json (same session, same reasoning:
#  paths/names/config, small, rewritten whole on every change - exactly
#  like the project registry).
#
#  Vector storage per KB: datas/<accountId>/kb/<kbId>.sqlite - handled
#  OUTSIDE the JSON registry encryption (too large/hot to decrypt+
#  re-encrypt whole on every search). Only the chunk TEXT column is
#  encrypted (AES-GCM, same per-account key); embeddings stay in
#  plaintext floats - raw numbers alone aren't meaningfully readable
#  without the chunk text next to them. This is the "(b)" option from
#  kiconnect-rag-spec.md section 2, not full SQLCipher - communicated
#  to the user as-is in the UI, no overstated security guarantee.
#
#  Embedding provider: deliberately generic OpenAI-compatible
#  /embeddings (baseUrl + model + optional apiKey) - NOT hardcoded to
#  any single local tool. Works with LM Studio, Ollama, vLLM,
#  text-generation-webui, AnythingLLM's built-in server, or any hosted
#  embeddings API that speaks this format (OpenAI, Mistral, DeepInfra,
#  Together, ...). Goes through the same is_allowed()/_pin_dns() SSRF
#  protection as the existing chat proxy, since the URL is user-supplied
#  and could point anywhere (localhost, LAN, or the public internet).
#
#  Vector search: plain SQLite + brute-force cosine similarity (numpy if
#  available, pure Python otherwise) rather than a sqlite-vec extension
#  dependency - a personal knowledge base's chunk count (typically
#  hundreds to low tens-of-thousands) makes a full scan fast enough
#  locally, and this way the feature works out of the box on every
#  platform without a native-extension build/install step that could
#  fail silently for some users. Swappable for sqlite-vec later
#  (Ausbaustufe 2/3) without changing the registry format or API shape -
#  only _kb_open_db()/_kb_cosine_search() would need to change.
# ════════════════════════════════════════════════════════════════
MAX_KB_FILE_SIZE = 25 * 1024 * 1024            # 25 MB/file - PDFs/DOCX are bigger than agent's 5MB text files
MAX_KB_FILES = 5000                            # safety cap per knowledge base
KB_IGNORE_DIRS = AGENT_IGNORE_DIRS             # reuse the agent's ignore list (.git, node_modules, ...)
KB_SUPPORTED_EXT = {
    '.txt', '.md', '.markdown', '.csv', '.tsv', '.log', '.json', '.yaml', '.yml',
    '.py', '.js', '.ts', '.jsx', '.tsx', '.html', '.css', '.c', '.cpp', '.h',
    '.java', '.go', '.rs', '.sh', '.sql', '.xml',
    '.pdf', '.docx', '.pptx', '.xlsx',
}
CHUNK_SIZE_TOKENS_DEFAULT = 512
CHUNK_OVERLAP_DEFAULT = 64
KB_TOP_K_DEFAULT = 8
KB_TOP_K_MAX = 30
_SAFE_KB_ID_RE = re.compile(r'^kb_[A-Za-z0-9]{1,64}$')

def _kb_clamp_top_k(value):
    try:
        return max(1, min(int(value), KB_TOP_K_MAX))
    except (TypeError, ValueError):
        return KB_TOP_K_DEFAULT


# ── Optional dependencies ─────────────────────────────────────────
# The KB feature degrades gracefully rather than crashing the whole
# proxy at import time if one of these isn't installed - each code path
# that needs one checks first and returns a clear error instead.
try:
    import sqlite3
except ImportError:
    sqlite3 = None
try:
    import pypdf as _pypdf
except ImportError:
    _pypdf = None
try:
    import docx as _python_docx
except ImportError:
    _python_docx = None
try:
    from pptx import Presentation as _PptxPresentation
except ImportError:
    _PptxPresentation = None
try:
    import openpyxl as _openpyxl
except ImportError:
    _openpyxl = None
try:
    import numpy as _np
except ImportError:
    _np = None
import struct as _struct
import hashlib as _hashlib

# ── Agent-session-backed registry (own key: 'kb_registry') ───────
def _load_kb_registry(sess):
    """Read+decrypt this account's knowledge-base registry. Same shape/
    pattern as _load_agent_registry() above."""
    fpath = _key_path(sess['accountId'], 'kb_registry')
    if not fpath or not os.path.isfile(fpath):
        return {'knowledgeBases': []}
    with _store_lock:
        with open(fpath, 'r', encoding='utf-8') as f:
            raw = f.read()
    try:
        blob = json.loads(raw) if raw else None
    except ValueError:
        blob = None
    if not blob:
        return {'knowledgeBases': []}
    data = _agent_decrypt(sess['key'], blob)
    if not isinstance(data.get('knowledgeBases'), list):
        data['knowledgeBases'] = []
    return data

def _save_kb_registry(sess, data):
    fpath = _key_path(sess['accountId'], 'kb_registry')
    if not fpath:
        return
    os.makedirs(os.path.dirname(fpath), exist_ok=True)
    blob = _agent_encrypt(sess['key'], data)
    with _store_lock:
        _atomic_write(fpath, json.dumps(blob).encode('utf-8'))

def _kb_validate_file_list(raw_paths):
    """Validates a list of absolute file paths for KB use (exists, is a
    file, supported extension, under the size cap, not inside a blocked
    root). Returns (accepted_realpaths, rejected_basenames) - used by both
    /kb/create (sourceType "files") and /kb/<id>/add-files, so "add
    individual files from various folders" behaves identically whichever
    endpoint it went through."""
    accepted, rejected = [], []
    seen = set()
    for raw in raw_paths[:MAX_KB_FILES]:
        raw = (raw or '').strip()
        if not raw:
            continue
        try:
            fpath = os.path.realpath(raw)
            # A folder source must not be a filesystem root, but an
            # individual file directly in a drive root is safe: the indexer
            # only receives that one file, not the complete drive.  Continue
            # protecting the application's own static/data directories.
            app_data_path = os.path.realpath(DATA_DIR)
            app_static_path = os.path.realpath(STATIC_DIR)
            in_app_storage = (
                fpath == app_data_path or fpath.startswith(app_data_path + os.sep)
                or fpath == app_static_path or fpath.startswith(app_static_path + os.sep)
            )
            ok = (
                os.path.isfile(fpath)
                and not in_app_storage
                and os.path.splitext(fpath)[1].lower() in KB_SUPPORTED_EXT
                and os.path.getsize(fpath) <= MAX_KB_FILE_SIZE
            )
        except OSError:
            ok = False
        if ok and fpath not in seen:
            accepted.append(fpath)
            seen.add(fpath)
        elif not ok:
            rejected.append(os.path.basename(raw) or raw)
    return accepted, rejected

def _kb_db_path(account_id, kb_id):
    adir = _account_dir(account_id)
    if not adir:
        return None
    kb_dir = os.path.join(adir, 'kb')
    os.makedirs(kb_dir, exist_ok=True)
    return os.path.join(kb_dir, kb_id + '.sqlite')

# ── Per-file text extraction ──────────────────────────────────────
def _kb_extract_pages(fpath, ext):
    """Returns (pages, error). `pages` is a list of (page_num_or_None,
    text) tuples - PDFs/PPTX get one entry per page/slide (so search
    results can cite "Seite 12"); everything else is a single entry with
    page=None. On failure, pages is None and error explains why."""
    ext = ext.lower()
    try:
        if ext == '.pdf':
            if not _pypdf:
                return None, 'pypdf not installed (pip install pypdf)'
            pages = []
            reader = _pypdf.PdfReader(fpath)
            for i, page in enumerate(reader.pages):
                try:
                    txt = page.extract_text() or ''
                except Exception:
                    txt = ''
                if txt.strip():
                    pages.append((i + 1, txt))
            return pages, None
        if ext == '.docx':
            if not _python_docx:
                return None, 'python-docx not installed (pip install python-docx)'
            doc = _python_docx.Document(fpath)
            txt = '\n'.join(p.text for p in doc.paragraphs if p.text)
            return ([(None, txt)] if txt.strip() else []), None
        if ext == '.pptx':
            if not _PptxPresentation:
                return None, 'python-pptx not installed (pip install python-pptx)'
            prs = _PptxPresentation(fpath)
            pages = []
            for i, slide in enumerate(prs.slides):
                parts = [shape.text for shape in slide.shapes if getattr(shape, 'has_text_frame', False) and shape.text]
                txt = '\n'.join(parts)
                if txt.strip():
                    pages.append((i + 1, txt))
            return pages, None
        if ext == '.xlsx':
            if not _openpyxl:
                return None, 'openpyxl not installed (pip install openpyxl)'
            # read_only+data_only: streams rows instead of loading the whole
            # workbook, and reads cached formula RESULTS rather than formula
            # strings (so "=SUM(A1:A9)" becomes "142", not the formula text).
            wb = _openpyxl.load_workbook(fpath, read_only=True, data_only=True)
            pages = []
            try:
                for i, ws in enumerate(wb.worksheets):
                    lines = []
                    for row in ws.iter_rows(values_only=True):
                        cells = ['' if c is None else str(c) for c in row]
                        if any(c.strip() for c in cells):
                            lines.append('\t'.join(cells))
                    txt = f'[{ws.title}]\n' + '\n'.join(lines)
                    if lines:
                        pages.append((i + 1, txt))
            finally:
                wb.close()
            return pages, None
        # Plain text / markdown / code / csv / etc. - read as UTF-8, best-effort.
        with open(fpath, 'r', encoding='utf-8', errors='ignore') as f:
            txt = f.read()
        return ([(None, txt)] if txt.strip() else []), None
    except Exception as e:
        return None, str(e)

# ── Chunking ───────────────────────────────────────────────────────
def _kb_chunk_text(text, chunk_tokens, overlap_tokens):
    """Word-based sliding-window chunking - word count as a simple,
    model-independent proxy for token count (see kiconnect-rag-spec.md
    5.2; exact token counting isn't critical for retrieval quality).
    Tries to end each chunk on a paragraph/sentence boundary within its
    last ~15% rather than a hard mid-sentence cut, when one is nearby."""
    words = text.split()
    if not words:
        return []
    chunk_tokens = max(32, int(chunk_tokens or CHUNK_SIZE_TOKENS_DEFAULT))
    overlap_tokens = max(0, min(int(overlap_tokens or 0), chunk_tokens - 1))
    step = max(1, chunk_tokens - overlap_tokens)
    n = len(words)
    chunks = []
    i = 0
    while i < n:
        end = min(i + chunk_tokens, n)
        chunk_text = ' '.join(words[i:end])
        if end < n:
            tail = chunk_text[-max(1, int(len(chunk_text) * 0.15)):]
            for marker in ('\n\n', '. ', '.\n', '! ', '? '):
                pos = tail.rfind(marker)
                if pos != -1:
                    cut = len(chunk_text) - len(tail) + pos + len(marker)
                    if cut > len(chunk_text) * 0.5:
                        chunk_text = chunk_text[:cut]
                    break
        chunk_text = chunk_text.strip()
        if chunk_text:
            chunks.append(chunk_text)
        if end >= n:
            break
        i += step
    return chunks

# ── Embedding call: generic OpenAI-compatible /embeddings ─────────
def _kb_get_embeddings(texts, embedding_cfg):
    """embedding_cfg: {baseUrl, model, apiKey?, lanConfirmed?}. Not tied
    to any single tool - works with anything speaking the OpenAI
    /embeddings shape (LM Studio, Ollama, vLLM, AnythingLLM, hosted
    APIs, ...). SSRF-protected the same way as the existing chat proxy
    (is_allowed()/_pin_dns() are defined further below in this file but
    resolved at call time, after the whole module has loaded)."""
    base_url = (embedding_cfg.get('baseUrl') or '').strip().rstrip('/')
    if not base_url:
        raise ValueError('No embedding server URL configured (see knowledge base settings).')
    if not texts:
        return []
    # Provider discovery uses /models and users sometimes paste the full
    # /embeddings URL from a curl example. Embeddings are a sibling endpoint,
    # so normalize either form back to the OpenAI-compatible base URL.
    base_url = re.sub(r'/(?:models|embeddings)$', '', base_url, flags=re.IGNORECASE)
    url = base_url + '/embeddings'
    ok, reason, needs_confirm, pinned_ip = is_allowed(url, 'POST', confirmed=bool(embedding_cfg.get('lanConfirmed')))
    if not ok:
        raise ValueError(reason + (' (confirm the address in the knowledge base settings first)' if needs_confirm else ''))
    headers = {'Content-Type': 'application/json'}
    api_key = (embedding_cfg.get('apiKey') or '').strip()
    headers['Authorization'] = f'Bearer {api_key or "not-needed"}'
    target_host = urlparse(url).hostname or ''
    with _pin_dns(target_host, pinned_ip):
        resp = requests.post(
            url, json={'model': embedding_cfg.get('model') or '', 'input': texts},
            headers=headers, timeout=(10, 120),
        )
    resp.raise_for_status()
    data = resp.json().get('data', [])
    data.sort(key=lambda d: d.get('index', 0))
    return [d['embedding'] for d in data]

# ── Vector store (plain SQLite, brute-force cosine search) ────────
_kb_locks = defaultdict(threading.Lock)   # kbId -> Lock, guards that KB's sqlite file
_kb_progress = {}                         # kbId -> {'status','done','total','error'} (in-RAM only)

def _kb_open_db(db_path):
    conn = sqlite3.connect(db_path)
    conn.execute('''CREATE TABLE IF NOT EXISTS chunks (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        source_file TEXT, page INTEGER, chunk_index INTEGER,
        text_encrypted TEXT, content_hash TEXT, embedding BLOB
    )''')
    conn.execute('CREATE INDEX IF NOT EXISTS idx_chunks_source ON chunks(source_file)')
    return conn

def _kb_pack_vec(vec):
    return _struct.pack(f'{len(vec)}f', *vec)

def _kb_unpack_vec(blob):
    n = len(blob) // 4
    return _struct.unpack(f'{n}f', blob)

def _kb_cosine_search(conn, query_vec, top_k):
    rows = conn.execute(
        'SELECT id, source_file, page, chunk_index, text_encrypted, embedding FROM chunks'
    ).fetchall()
    if not rows:
        return []
    scored = []
    if _np is not None:
        q = _np.asarray(query_vec, dtype='float32')
        qn = q / (float(_np.linalg.norm(q)) or 1e-9)
        for rid, src, page, cidx, text_enc, emb_blob in rows:
            v = _np.frombuffer(emb_blob, dtype='float32')
            vn = v / (float(_np.linalg.norm(v)) or 1e-9)
            scored.append((float(_np.dot(qn, vn)), rid, src, page, cidx, text_enc))
    else:
        import math
        qnorm = math.sqrt(sum(x * x for x in query_vec)) or 1e-9
        for rid, src, page, cidx, text_enc, emb_blob in rows:
            v = _kb_unpack_vec(emb_blob)
            dot = sum(a * b for a, b in zip(query_vec, v))
            vnorm = math.sqrt(sum(x * x for x in v)) or 1e-9
            scored.append((dot / (qnorm * vnorm), rid, src, page, cidx, text_enc))
    scored.sort(key=lambda t: t[0], reverse=True)
    return scored[:top_k]

# ── Background indexing ────────────────────────────────────────────
def _kb_run_index(account_id, kb_id, sess_key, full_reindex=True):
    """Runs in a background thread (started from /kb/create and
    /kb/reindex) - walks the KB's folder, extracts+chunks+embeds+stores
    every supported file, updating _kb_progress as it goes so
    /kb/<id>/status can report it without blocking the chat. Writes the
    final result into the persistent (encrypted) registry when done."""
    progress = _kb_progress.setdefault(kb_id, {})
    progress.update({'status': 'indexing', 'done': 0, 'total': 0, 'error': None})
    sess = {'accountId': account_id, 'key': sess_key, 'ts': time.time()}
    try:
        registry = _load_kb_registry(sess)
        entry = next((k for k in registry['knowledgeBases'] if k.get('id') == kb_id), None)
        if not entry:
            progress.update({'status': 'error', 'error': 'Knowledge base was deleted.'})
            return
        root = os.path.realpath(entry.get('path') or '') if entry.get('path') else ''
        if root and (not os.path.isdir(root) or _is_blocked_root(root)):
            progress.update({'status': 'error', 'error': 'Folder no longer accessible.'})
            return
        settings = entry.get('settings') or {}
        embedding_cfg = settings.get('embedding') or {}
        chunk_tokens = settings.get('chunkTokens') or CHUNK_SIZE_TOKENS_DEFAULT
        chunk_overlap = settings.get('chunkOverlap') or CHUNK_OVERLAP_DEFAULT

        # `files` maps every indexed path to the label used for citations
        # and dedup ("source"): for the folder root, that's the path
        # relative to the root (so "sub/dir/file.md"); for individually
        # added files (any folder, including ones added on top of a
        # folder-based KB via /kb/<id>/add-files) it's "parentDir/file.ext"
        # so files with the same name from different folders stay distinct.
        files = []      # list of (abs_path, label)
        seen_abs = set()
        excluded_abs = {os.path.realpath(p) for p in (entry.get('excludedFiles') or [])}
        if root:
            for dirpath, dirs, filenames in os.walk(root):
                dirs[:] = sorted(d for d in dirs if d not in KB_IGNORE_DIRS and not d.startswith('.'))
                for fn in sorted(filenames):
                    if fn.startswith('.') or os.path.splitext(fn)[1].lower() not in KB_SUPPORTED_EXT:
                        continue
                    fpath = os.path.join(dirpath, fn)
                    try:
                        if os.path.realpath(fpath) in excluded_abs or os.path.getsize(fpath) > MAX_KB_FILE_SIZE:
                            continue
                    except OSError:
                        continue
                    rel = os.path.relpath(fpath, root).replace('\\', '/')
                    files.append((fpath, rel))
                    seen_abs.add(os.path.realpath(fpath))
                    if len(files) >= MAX_KB_FILES:
                        break
                if len(files) >= MAX_KB_FILES:
                    break
        for fpath in (entry.get('fileList') or []):
            if len(files) >= MAX_KB_FILES:
                break
            try:
                fpath = os.path.realpath(fpath)
                if not os.path.isfile(fpath) or fpath in seen_abs or fpath in excluded_abs:
                    continue
                if os.path.getsize(fpath) > MAX_KB_FILE_SIZE:
                    continue
            except OSError:
                continue
            parent_name = os.path.basename(os.path.dirname(fpath))
            label = f'{parent_name}/{os.path.basename(fpath)}' if parent_name else os.path.basename(fpath)
            files.append((fpath, label))
            seen_abs.add(fpath)

        db_path = _kb_db_path(account_id, kb_id)
        lock = _kb_locks[kb_id]
        if full_reindex:
            with lock:
                conn = _kb_open_db(db_path)
                conn.execute('DELETE FROM chunks')
                conn.commit()
                conn.close()
        else:
            # Incremental run (e.g. /kb/<id>/add-files): only embed files not
            # already represented in the index - v1 has no change detection
            # (see kb_reindex()'s comment), so a file already indexed once is
            # assumed unchanged and is skipped here rather than duplicated.
            with lock:
                conn = _kb_open_db(db_path)
                already = {row[0] for row in conn.execute('SELECT DISTINCT source_file FROM chunks')}
                conn.close()
            files = [(fpath, rel) for fpath, rel in files if rel not in already]
        progress['total'] = len(files)

        failed_files = []
        chunk_count = 0
        embed_dim = None
        for fpath, rel in files:
            ext = os.path.splitext(fpath)[1].lower()
            pages, err = _kb_extract_pages(fpath, ext)
            if err or not pages:
                if err:
                    failed_files.append(f'{rel} ({err})')
                progress['done'] += 1
                continue

            to_embed = []  # (page_num, chunk_idx, text)
            for page_num, page_text in pages:
                for idx, piece in enumerate(_kb_chunk_text(page_text, chunk_tokens, chunk_overlap)):
                    to_embed.append((page_num, idx, piece))
            if not to_embed:
                progress['done'] += 1
                continue

            try:
                BATCH = 64
                all_vecs = []
                for b in range(0, len(to_embed), BATCH):
                    batch_texts = [t[2] for t in to_embed[b:b + BATCH]]
                    all_vecs.extend(_kb_get_embeddings(batch_texts, embedding_cfg))
                if all_vecs and embed_dim is None:
                    embed_dim = len(all_vecs[0])
            except Exception as e:
                failed_files.append(f'{rel} ({e})')
                progress['done'] += 1
                continue

            with lock:
                conn = _kb_open_db(db_path)
                for (page_num, idx, piece), vec in zip(to_embed, all_vecs):
                    text_enc = _agent_encrypt(sess_key, {'t': piece})
                    conn.execute(
                        'INSERT INTO chunks (source_file, page, chunk_index, text_encrypted, content_hash, embedding) '
                        'VALUES (?,?,?,?,?,?)',
                        (rel, page_num, idx, text_enc, _hashlib.sha256(piece.encode('utf-8')).hexdigest(), _kb_pack_vec(vec))
                    )
                    chunk_count += 1
                conn.commit()
                conn.close()
            progress['done'] += 1

        with _store_lock:
            registry = _load_kb_registry(sess)
            entry = next((k for k in registry['knowledgeBases'] if k.get('id') == kb_id), None)
            if entry:
                entry['indexState'] = {
                    'status': 'ready' if not failed_files else 'ready_with_errors',
                    'fileCount': len(files) - len(failed_files),
                    'failedFiles': failed_files[:50],
                    'chunkCount': chunk_count,
                    'embeddingDim': embed_dim,
                    'lastIndexed': time.time(),
                }
                _save_kb_registry(sess, registry)
        progress['status'] = 'ready'
    except Exception as e:
        progress.update({'status': 'error', 'error': str(e)})

def _kb_error(msg, status=400):
    return Response(json.dumps({'error': msg}), status, content_type='application/json')

def _kb_ok(payload=None, status=200):
    return Response(json.dumps(payload if payload is not None else {'ok': True}), status, content_type='application/json')

# Every /kb/* view used to repeat the same five lines by hand (OPTIONS
# preflight, session check, kb_id format check). This decorator does it
# once: it handles OPTIONS, requires a valid agent session and passes it
# in as the view's first argument, and - if the route has a <kb_id> - checks
# its format before the view ever sees it.
def _kb_route(f):
    @wraps(f)
    def wrapper(*args, **kwargs):
        if request.method == 'OPTIONS':
            return Response('', 204, headers={**CORS_HEADERS, **SECURITY_HEADERS})
        sess = _agent_session_or_401()
        if not sess:
            return _kb_error('Session expired - please log in again.', 401)
        kb_id = kwargs.get('kb_id')
        if kb_id is not None and not _SAFE_KB_ID_RE.match(kb_id):
            return _kb_error('Invalid id.', 400)
        return f(sess, *args, **kwargs)
    return wrapper

# Looks up a knowledge base by id in an already-loaded registry. Returns
# (entry, None) on success or (None, error_response) - callers just do
# `entry, err = _kb_require_entry(...); if err: return err`, replacing the
# four-line "find or 404" block every entry-scoped route used to repeat.
def _kb_require_entry(registry, kb_id):
    entry = next((k for k in registry['knowledgeBases'] if k.get('id') == kb_id), None)
    if not entry:
        return None, _kb_error('Knowledge base not found.', 404)
    return entry, None

# ── /kb/create - register a folder or an explicit file list as a new
#    knowledge base and start indexing ──
@app.route('/kb/create', methods=['POST', 'OPTIONS'])
@_kb_route
def kb_create(sess):
    if sqlite3 is None:
        return _kb_error('sqlite3 is not available in this Python installation.', 500)
    try: body = request.get_json(force=True, silent=True) or {}
    except Exception: body = {}
    name = (body.get('name') or '').strip()
    source_type = (body.get('sourceType') or 'folder').strip()
    if not name or len(name) > 96:
        return _kb_error('Invalid name (1-96 characters).')
    if source_type not in ('folder', 'files'):
        return _kb_error('sourceType must be "folder" or "files".')

    target = ''
    file_list = []
    if source_type == 'folder':
        raw_path = (body.get('path') or '').strip()
        if not raw_path:
            return _kb_error('Please provide a folder path.')
        target = os.path.realpath(raw_path)
        if not os.path.isdir(target):
            return _kb_error('Folder not found.', 404)
        if _is_blocked_root(target):
            return _kb_error('This folder is not allowed for security reasons.', 403)
    else:
        # "Individual files" mode: any number of files, from any number of
        # folders - no single common root required (kiconnect-rag-spec.md
        # section 9 Phase 2). Silently skips paths that don't validate
        # rather than failing the whole create; the skipped ones are
        # reported back in `filesRejected` so the UI can show them.
        raw_list = body.get('fileList') or []
        allow_empty_upload = bool(body.get('allowEmptyFileList'))
        if not isinstance(raw_list, list) or (not raw_list and not allow_empty_upload):
            return _kb_error('Please provide at least one file.')
        file_list, rejected = _kb_validate_file_list(raw_list)
        if not file_list and not allow_empty_upload:
            return _kb_error('None of the given files could be used.')

    embedding = body.get('embedding') or {}
    with _store_lock:
        registry = _load_kb_registry(sess)
        kb_id = 'kb_' + os.urandom(8).hex()
        entry = {
            'id': kb_id, 'name': name, 'sourceType': source_type,
            'path': target, 'fileList': file_list,
            'createdAt': time.time(),
            'settings': {
                'embedding': {
                    'baseUrl': (embedding.get('baseUrl') or '').strip(),
                    'model': (embedding.get('model') or '').strip(),
                    'apiKey': (embedding.get('apiKey') or '').strip(),
                },
                'chunkTokens': int(body.get('chunkTokens') or CHUNK_SIZE_TOKENS_DEFAULT),
                'chunkOverlap': int(body.get('chunkOverlap') or CHUNK_OVERLAP_DEFAULT),
                'reranker': bool(body.get('reranker')),
                'topK': _kb_clamp_top_k(body.get('topK') or KB_TOP_K_DEFAULT),
            },
            'indexState': {
                'status': 'awaiting_upload' if source_type == 'files' and not file_list else 'pending',
                'fileCount': 0, 'chunkCount': 0, 'lastIndexed': None,
            },
        }
        registry['knowledgeBases'].append(entry)
        _save_kb_registry(sess, registry)
    # A drag-and-drop-only KB is intentionally created empty and receives
    # its browser files in the following /upload-files request.  Starting an
    # empty index job here races that request and makes it fail with 409.
    if source_type == 'folder' or file_list:
        threading.Thread(target=_kb_run_index, args=(sess['accountId'], kb_id, sess['key'], True), daemon=True).start()
    resp = {'id': kb_id, 'name': name, 'path': target}
    if source_type == 'files':
        resp['filesAdded'] = len(file_list)
        if rejected:
            resp['filesRejected'] = rejected
    return _kb_ok(resp)

# ── /kb/list ──
@app.route('/kb/list', methods=['GET', 'OPTIONS'])
@_kb_route
def kb_list(sess):
    registry = _load_kb_registry(sess)
    out = []
    for kb in registry['knowledgeBases']:
        raw_path = kb.get('path') or ''
        target = os.path.realpath(raw_path) if raw_path else ''
        file_list = kb.get('fileList') or []
        folder_missing = bool(target) and not (os.path.isdir(target) and not _is_blocked_root(target))
        existing_files = [f for f in file_list if os.path.isfile(f)]
        # "missing" means nothing this KB points at is reachable any more -
        # a folder-based KB with a few extra individually-added files that
        # went away isn't "missing", only fully empty/unreachable ones are.
        missing = folder_missing and not existing_files if target else (bool(file_list) and not existing_files)
        settings = dict(kb.get('settings') or {})
        emb = dict(settings.get('embedding') or {})
        if emb.get('apiKey'):
            emb['apiKey'] = '••••••••'   # never send the real key back in a listing
        settings['embedding'] = emb
        folders = {os.path.dirname(f) for f in file_list}
        out.append({
            'id': kb['id'], 'name': kb['name'], 'sourceType': kb.get('sourceType'),
            'path': target, 'fileCount': len(file_list), 'folderCount': len(folders),
            'missing': missing,
            'settings': settings, 'indexState': kb.get('indexState') or {},
        })
    return _kb_ok({'knowledgeBases': out})

# ── /kb/<id>/status - polling target while indexing runs in the background ──
@app.route('/kb/<kb_id>/status', methods=['GET', 'OPTIONS'])
@_kb_route
def kb_status(sess, kb_id):
    live = _kb_progress.get(kb_id)
    if live:
        return _kb_ok({'status': live.get('status'), 'done': live.get('done', 0),
                       'total': live.get('total', 0), 'error': live.get('error')})
    registry = _load_kb_registry(sess)
    entry, err = _kb_require_entry(registry, kb_id)
    if err: return err
    state = entry.get('indexState') or {}
    fc = state.get('fileCount', 0)
    return _kb_ok({'status': state.get('status', 'pending'), 'done': fc, 'total': fc, 'error': None})

# ── /kb/<id>/sources - list and remove individual indexed sources ──────
@app.route('/kb/<kb_id>/sources', methods=['GET', 'DELETE', 'OPTIONS'])
@_kb_route
def kb_sources(sess, kb_id):
    if _kb_progress.get(kb_id, {}).get('status') in ('pending', 'indexing'):
        return _kb_error('Please wait for indexing to finish first.', 409)
    with _store_lock:
        registry = _load_kb_registry(sess)
        entry, err = _kb_require_entry(registry, kb_id)
        if err: return err
        root = os.path.realpath(entry.get('path') or '') if entry.get('path') else ''
        excluded = {os.path.realpath(p) for p in (entry.get('excludedFiles') or [])}
        sources, seen = [], set()
        def add_source(path, label, source_type):
            path = os.path.realpath(path)
            if path in seen or path in excluded or not os.path.isfile(path): return
            seen.add(path); sources.append({'path': path, 'label': label, 'sourceType': source_type})
        if root and os.path.isdir(root):
            for dirpath, dirs, filenames in os.walk(root):
                dirs[:] = [d for d in dirs if d not in KB_IGNORE_DIRS and not d.startswith('.')]
                for fn in filenames:
                    path = os.path.join(dirpath, fn)
                    if not fn.startswith('.') and os.path.splitext(fn)[1].lower() in KB_SUPPORTED_EXT:
                        add_source(path, os.path.relpath(path, root).replace('\\', '/'), 'folder')
        for path in entry.get('fileList') or []:
            path = os.path.realpath(path)
            parent = os.path.basename(os.path.dirname(path))
            add_source(path, f'{parent}/{os.path.basename(path)}' if parent else os.path.basename(path), 'file')
        if request.method == 'GET':
            return _kb_ok({'sources': sources})
        try: body = request.get_json(force=True, silent=True) or {}
        except Exception: body = {}
        paths = body.get('paths') or []
        if not isinstance(paths, list) or not paths:
            return _kb_error('Please provide at least one source path.')
        selected = {os.path.realpath(p) for p in paths}
        removable = [s for s in sources if s['path'] in selected]
        if not removable:
            return _kb_error('No matching sources found.', 404)
        labels = [s['label'] for s in removable]
        entry['fileList'] = [p for p in (entry.get('fileList') or []) if os.path.realpath(p) not in selected]
        entry['excludedFiles'] = sorted(excluded | {s['path'] for s in removable if s['sourceType'] == 'folder'})
        _save_kb_registry(sess, registry)
    db_path = _kb_db_path(sess['accountId'], kb_id)
    file_count, chunk_count = 0, 0
    if db_path and os.path.isfile(db_path):
        with _kb_locks[kb_id]:
            conn = _kb_open_db(db_path)
            conn.executemany('DELETE FROM chunks WHERE source_file = ?', [(label,) for label in labels])
            file_count = conn.execute('SELECT COUNT(DISTINCT source_file) FROM chunks').fetchone()[0]
            chunk_count = conn.execute('SELECT COUNT(*) FROM chunks').fetchone()[0]
            conn.commit(); conn.close()
    with _store_lock:
        registry = _load_kb_registry(sess)
        entry = next((k for k in registry['knowledgeBases'] if k.get('id') == kb_id), None)
        if entry:
            state = entry.setdefault('indexState', {})
            state['fileCount'], state['chunkCount'] = file_count, chunk_count
            _save_kb_registry(sess, registry)
    return _kb_ok({'removed': len(removable)})

# ── /kb/<id>/reindex - full re-index (no change detection in v1, see spec 9/Phase2) ──
@app.route('/kb/<kb_id>/reindex', methods=['POST', 'OPTIONS'])
@_kb_route
def kb_reindex(sess, kb_id):
    with _store_lock:
        registry = _load_kb_registry(sess)
        entry, err = _kb_require_entry(registry, kb_id)
        if err: return err
        if _kb_progress.get(kb_id, {}).get('status') in ('pending', 'indexing'):
            return _kb_error('Already indexing.', 409)
        # Persist the transition before starting the background thread. The
        # list endpoint drives the browser's progress polling, so leaving a
        # stale "ready" state here made re-indexing appear to do nothing.
        entry['indexState'] = {'status': 'pending', 'fileCount': 0, 'chunkCount': 0, 'lastIndexed': None}
        _save_kb_registry(sess, registry)
        _kb_progress[kb_id] = {'status': 'pending', 'done': 0, 'total': 0, 'error': None}
    threading.Thread(target=_kb_run_index, args=(sess['accountId'], kb_id, sess['key'], True), daemon=True).start()
    return _kb_ok({'ok': True})

# ── GET /kb/<id>/export - decrypt every chunk and hand back a plain-JSON
#    dump the browser turns into a download. Deliberately NOT encrypted:
#    the whole point is portability to a *different* account/password, and
#    that account's key can't decrypt bytes sealed with this one's. The
#    frontend must show a clear "this file is unencrypted, store it
#    somewhere safe" notice before/at the actual download (kiconnect.js).
#    The embedding apiKey is never included - it's a per-account credential,
#    not knowledge-base data, and has no business leaving the account. ──
EXPORT_FORMAT_VERSION = 1

@app.route('/kb/<kb_id>/export', methods=['GET', 'OPTIONS'])
@_kb_route
def kb_export(sess, kb_id):
    registry = _load_kb_registry(sess)
    entry, err = _kb_require_entry(registry, kb_id)
    if err: return err
    db_path = _kb_db_path(sess['accountId'], kb_id)
    chunks_out = []
    if db_path and os.path.isfile(db_path):
        lock = _kb_locks[kb_id]
        with lock:
            conn = _kb_open_db(db_path)
            try:
                rows = conn.execute(
                    'SELECT source_file, page, chunk_index, text_encrypted, content_hash, embedding FROM chunks'
                ).fetchall()
            finally:
                conn.close()
        for src, page, cidx, text_enc, chash, emb_blob in rows:
            try:
                text = _agent_decrypt(sess['key'], text_enc).get('t', '')
            except Exception:
                continue  # skip a chunk we can't decrypt rather than fail the whole export
            chunks_out.append({
                'sourceFile': os.path.basename(src or ''),  # basename only - the exporter's local
                'page': page,                                # folder layout isn't the importer's business
                'chunkIndex': cidx,
                'text': text,
                'contentHash': chash,
                'embedding': list(_kb_unpack_vec(emb_blob)) if emb_blob else None,
            })
    settings = entry.get('settings') or {}
    embedding_cfg = settings.get('embedding') or {}
    return _kb_ok({
        'kiconnectExport': 'kb/v1',
        'formatVersion': EXPORT_FORMAT_VERSION,
        'exportedAt': time.time(),
        'name': entry.get('name') or 'Knowledge base',
        'settings': {
            # baseUrl/model travel along as a hint for the importing account's
            # own embedding setup - the apiKey does not (see comment above).
            'embedding': {'baseUrl': embedding_cfg.get('baseUrl') or '', 'model': embedding_cfg.get('model') or ''},
            'chunkTokens': settings.get('chunkTokens') or CHUNK_SIZE_TOKENS_DEFAULT,
            'chunkOverlap': settings.get('chunkOverlap') or CHUNK_OVERLAP_DEFAULT,
            'reranker': bool(settings.get('reranker')),
            'topK': _kb_clamp_top_k(settings.get('topK') or KB_TOP_K_DEFAULT),
        },
        'chunkCount': len(chunks_out),
        'chunks': chunks_out,
    })

# ── POST /kb/import - opposite direction: takes a kb/v1 export (this
#    account's password, whatever it is, was never involved in producing
#    that plaintext) and re-encrypts every chunk under *this* session's key
#    before it ever touches disk. Existing embeddings are reused as-is
#    (no re-embedding call, no embedding server needed) provided the
#    importing account intends to search with the same embedding model -
#    mixing models in one KB would make cosine scores meaningless, so a
#    mismatched model name is flagged as a warning, not a hard error, since
#    we can't verify dimensions until something is actually indexed. ──
@app.route('/kb/import', methods=['POST', 'OPTIONS'])
@_kb_route
def kb_import(sess):
    if sqlite3 is None:
        return _kb_error('sqlite3 is not available in this Python installation.', 500)
    try: body = request.get_json(force=True, silent=True) or {}
    except Exception: body = {}
    export = body.get('export') or {}
    if export.get('kiconnectExport') != 'kb/v1':
        return _kb_error('This file is not a recognized knowledge-base export.')
    chunks_in = export.get('chunks')
    if not isinstance(chunks_in, list):
        return _kb_error('Export file has no chunks array.')
    if len(chunks_in) > 200000:
        return _kb_error('Export is too large to import in one request.')

    name = (body.get('name') or export.get('name') or 'Imported knowledge base').strip()[:96]
    src_settings = export.get('settings') or {}
    embedding_cfg = src_settings.get('embedding') or {}
    # An apiKey can optionally be supplied fresh in the import request itself
    # (the export never carries one) so the KB is immediately searchable.
    new_embedding = body.get('embedding') or {}

    with _store_lock:
        registry = _load_kb_registry(sess)
        kb_id = 'kb_' + os.urandom(8).hex()
        entry = {
            'id': kb_id, 'name': name, 'sourceType': 'import',
            'path': '', 'fileList': [],
            'createdAt': time.time(),
            'settings': {
                'embedding': {
                    'baseUrl': (new_embedding.get('baseUrl') or embedding_cfg.get('baseUrl') or '').strip(),
                    'model': (new_embedding.get('model') or embedding_cfg.get('model') or '').strip(),
                    'apiKey': (new_embedding.get('apiKey') or '').strip(),
                },
                'chunkTokens': int(src_settings.get('chunkTokens') or CHUNK_SIZE_TOKENS_DEFAULT),
                'chunkOverlap': int(src_settings.get('chunkOverlap') or CHUNK_OVERLAP_DEFAULT),
                'reranker': bool(src_settings.get('reranker')),
                'topK': _kb_clamp_top_k(src_settings.get('topK') or KB_TOP_K_DEFAULT),
            },
            'indexState': {'status': 'indexing', 'fileCount': 0, 'chunkCount': 0, 'lastIndexed': None},
        }
        registry['knowledgeBases'].append(entry)
        _save_kb_registry(sess, registry)

    db_path = _kb_db_path(sess['accountId'], kb_id)
    if not db_path:
        return _kb_error('Unable to prepare storage for the imported knowledge base.', 500)
    MAX_IMPORT_CHUNK_CHARS = 20000  # guards against a crafted export ballooning memory/disk
    embed_dim, inserted, sources, skipped = None, 0, set(), 0
    conn = _kb_open_db(db_path)
    try:
        for c in chunks_in:
            text = (c.get('text') or '').strip()[:MAX_IMPORT_CHUNK_CHARS]
            emb = c.get('embedding')
            if not text or not isinstance(emb, list) or not emb:
                skipped += 1
                continue
            try:
                vec = [float(x) for x in emb]
            except (TypeError, ValueError):
                skipped += 1
                continue
            if embed_dim is None:
                embed_dim = len(vec)
            elif len(vec) != embed_dim:
                skipped += 1  # inconsistent vector size - can't sit in the same brute-force index
                continue
            src = os.path.basename(c.get('sourceFile') or 'imported')
            sources.add(src)
            text_enc = _agent_encrypt(sess['key'], {'t': text})
            content_hash = _hashlib.sha256(text.encode('utf-8')).hexdigest()
            conn.execute(
                'INSERT INTO chunks (source_file, page, chunk_index, text_encrypted, content_hash, embedding) '
                'VALUES (?, ?, ?, ?, ?, ?)',
                (src, c.get('page'), c.get('chunkIndex') or 0, text_enc, content_hash, _kb_pack_vec(vec))
            )
            inserted += 1
        conn.commit()
    finally:
        conn.close()

    with _store_lock:
        registry = _load_kb_registry(sess)
        entry = next((k for k in registry['knowledgeBases'] if k.get('id') == kb_id), None)
        if entry:
            entry['fileList'] = []
            entry['indexState'] = {
                'status': 'ready', 'fileCount': len(sources), 'chunkCount': inserted,
                'lastIndexed': time.time(), 'embeddingDim': embed_dim,
            }
            _save_kb_registry(sess, registry)

    resp = {'id': kb_id, 'name': name, 'imported': inserted, 'skipped': skipped}
    if embedding_cfg.get('model') and embedding_cfg.get('model') != entry['settings']['embedding']['model']:
        resp['warning'] = 'Embedding model differs from the source KB - re-index for consistent search results.'
    return _kb_ok(resp)

# ── DELETE /kb/<id> - unregister + remove its vector DB (source folder untouched) ──
@app.route('/kb/<kb_id>', methods=['DELETE', 'OPTIONS'])
@_kb_route
def kb_delete(sess, kb_id):
    with _store_lock:
        registry = _load_kb_registry(sess)
        before = len(registry['knowledgeBases'])
        registry['knowledgeBases'] = [k for k in registry['knowledgeBases'] if k.get('id') != kb_id]
        if len(registry['knowledgeBases']) == before:
            return _kb_error('Knowledge base not found.', 404)
        _save_kb_registry(sess, registry)
    _kb_progress.pop(kb_id, None)
    db_path = _kb_db_path(sess['accountId'], kb_id)
    try:
        if db_path and os.path.isfile(db_path):
            os.remove(db_path)
    except OSError:
        pass
    return _kb_ok({'ok': True})

# ── /kb/<id>/search - embed the query, top-k cosine search, decrypt+return chunks ──
@app.route('/kb/<kb_id>/search', methods=['POST', 'OPTIONS'])
@_kb_route
def kb_search(sess, kb_id):
    try: body = request.get_json(force=True, silent=True) or {}
    except Exception: body = {}
    query = (body.get('query') or '').strip()
    if not query:
        return _kb_error('No query provided.')
    registry = _load_kb_registry(sess)
    entry, err = _kb_require_entry(registry, kb_id)
    if err: return err
    # An explicit topK in the request wins (e.g. the composer merging several
    # active KBs); otherwise fall back to this KB's own configured topK
    # (Settings ▸ this knowledge base ▸ Advanced), not a single hardcoded value.
    stored_top_k = (entry.get('settings') or {}).get('topK', KB_TOP_K_DEFAULT)
    top_k = _kb_clamp_top_k(body.get('topK')) if 'topK' in body else _kb_clamp_top_k(stored_top_k)
    embedding_cfg = (entry.get('settings') or {}).get('embedding') or {}
    db_path = _kb_db_path(sess['accountId'], kb_id)
    if not db_path or not os.path.isfile(db_path):
        return _kb_ok({'results': []})
    try:
        vecs = _kb_get_embeddings([query], embedding_cfg)
    except Exception as e:
        return _kb_error(f'Embedding request failed: {e}', 502)
    if not vecs:
        return _kb_ok({'results': []})
    lock = _kb_locks[kb_id]
    with lock:
        conn = _kb_open_db(db_path)
        try:
            scored = _kb_cosine_search(conn, vecs[0], top_k)
        finally:
            conn.close()
    results = []
    for score, rid, src, page, cidx, text_enc in scored:
        try:
            text = _agent_decrypt(sess['key'], text_enc).get('t', '')
        except Exception:
            text = ''
        results.append({'text': text, 'source': src, 'page': page, 'chunkIndex': cidx, 'score': round(score, 4)})
    return _kb_ok({'results': results})

# ── PATCH /kb/<id>/settings - name/embedding provider/chunking/reranker/topK ──
@app.route('/kb/<kb_id>/settings', methods=['PATCH', 'OPTIONS'])
@_kb_route
def kb_update_settings(sess, kb_id):
    try: body = request.get_json(force=True, silent=True) or {}
    except Exception: body = {}
    with _store_lock:
        registry = _load_kb_registry(sess)
        entry, err = _kb_require_entry(registry, kb_id)
        if err: return err
        settings = entry.setdefault('settings', {})
        if 'embedding' in body:
            emb = settings.setdefault('embedding', {})
            new_emb = body.get('embedding') or {}
            if 'baseUrl' in new_emb: emb['baseUrl'] = (new_emb.get('baseUrl') or '').strip()
            if 'model' in new_emb: emb['model'] = (new_emb.get('model') or '').strip()
            # Only overwrite the stored key if a real (non-masked) value was sent -
            # kb_list() only ever returns "••••••••" as a placeholder, never the key itself.
            if 'apiKey' in new_emb and new_emb.get('apiKey') != '••••••••':
                emb['apiKey'] = (new_emb.get('apiKey') or '').strip()
        if 'chunkTokens' in body:
            settings['chunkTokens'] = int(body.get('chunkTokens') or CHUNK_SIZE_TOKENS_DEFAULT)
        if 'chunkOverlap' in body:
            settings['chunkOverlap'] = int(body.get('chunkOverlap') or CHUNK_OVERLAP_DEFAULT)
        if 'reranker' in body:
            settings['reranker'] = bool(body.get('reranker'))
        if 'topK' in body:
            settings['topK'] = _kb_clamp_top_k(body.get('topK'))
        if (body.get('name') or '').strip():
            entry['name'] = body['name'].strip()[:96]
        _save_kb_registry(sess, registry)
    return _kb_ok({'ok': True})

# ── POST /kb/<id>/add-files - append individual files (from any folder(s))
#    to an existing knowledge base, works for both "folder" and "files"
#    KBs. Only embeds+indexes the newly added files (full_reindex=False),
#    not a full re-index. (kiconnect-rag-spec.md section 9 Phase 2) ──────
@app.route('/kb/<kb_id>/add-files', methods=['POST', 'OPTIONS'])
@_kb_route
def kb_add_files(sess, kb_id):
    try: body = request.get_json(force=True, silent=True) or {}
    except Exception: body = {}
    raw_paths = body.get('paths') or []
    if not isinstance(raw_paths, list) or not raw_paths:
        return _kb_error('Please provide at least one file.')
    if _kb_progress.get(kb_id, {}).get('status') == 'indexing':
        return _kb_error('Already indexing - please wait for it to finish first.', 409)
    with _store_lock:
        registry = _load_kb_registry(sess)
        entry, err = _kb_require_entry(registry, kb_id)
        if err: return err
        existing = {os.path.realpath(p) for p in (entry.get('fileList') or [])}
        accepted, rejected = _kb_validate_file_list(raw_paths)
        new_files = [p for p in accepted if p not in existing]
        already = [os.path.basename(p) for p in accepted if p in existing]
        if not new_files and not rejected and not already:
            return _kb_error('No usable files given.')
        entry.setdefault('fileList', [])
        entry['fileList'].extend(new_files)
        if new_files:
            entry['indexState'] = {'status': 'pending', 'fileCount': 0, 'chunkCount': 0, 'lastIndexed': None}
        _save_kb_registry(sess, registry)
    if new_files:
        # Incremental: re-embeds/inserts only the newly added files, leaves
        # the rest of the index untouched (full_reindex=False keeps the
        # existing chunk rows instead of wiping the whole KB first).
        _kb_progress[kb_id] = {'status': 'pending', 'done': 0, 'total': 0, 'error': None}
        threading.Thread(target=_kb_run_index, args=(sess['accountId'], kb_id, sess['key'], False), daemon=True).start()
    resp = {'added': len(new_files)}
    if rejected or already:
        resp['rejected'] = rejected + already
    return _kb_ok(resp)


# ── POST /kb/<id>/upload-files - receive browser drag & drop files ──────
@app.route('/kb/<kb_id>/upload-files', methods=['POST', 'OPTIONS'])
@_kb_route
def kb_upload_files(sess, kb_id):
    try: body = request.get_json(force=True, silent=True) or {}
    except Exception: body = {}
    files = body.get('files') or []
    if not isinstance(files, list) or not files:
        return _kb_error('Please provide at least one file.')
    if len(files) > MAX_KB_FILES:
        return _kb_error(f'Too many files (maximum {MAX_KB_FILES}).')
    if _kb_progress.get(kb_id, {}).get('status') == 'indexing':
        return _kb_error('Already indexing - please wait for it to finish first.', 409)

    # Validate ownership before accepting bytes, so an invalid id cannot
    # leave unreferenced uploaded files behind.
    registry = _load_kb_registry(sess)
    if not any(k.get('id') == kb_id for k in registry['knowledgeBases']):
        return _kb_error('Knowledge base not found.', 404)

    upload_dir = os.path.join(_account_dir(sess['accountId']) or '', 'kb_uploads', kb_id)
    if not upload_dir or not os.path.realpath(upload_dir).startswith(os.path.realpath(_account_dir(sess['accountId']) or '') + os.sep):
        return _kb_error('Unable to prepare upload storage.', 500)
    os.makedirs(upload_dir, exist_ok=True)
    stored, rejected = [], []
    for item in files:
        name = os.path.basename(str((item or {}).get('name') or ''))
        ext = os.path.splitext(name)[1].lower()
        encoded = (item or {}).get('dataBase64') or ''
        if not name or ext not in KB_SUPPORTED_EXT:
            rejected.append(name or 'unnamed file')
            continue
        try:
            data = base64.b64decode(encoded, validate=True)
        except Exception:
            rejected.append(name)
            continue
        if not data or len(data) > MAX_KB_FILE_SIZE:
            rejected.append(name)
            continue
        safe_name = re.sub(r'[^A-Za-z0-9._ -]', '_', name).strip(' .') or ('upload' + ext)
        target = os.path.join(upload_dir, os.urandom(6).hex() + '_' + safe_name)
        try:
            with open(target, 'wb') as f:
                f.write(data)
            stored.append(target)
        except OSError:
            rejected.append(name)

    with _store_lock:
        registry = _load_kb_registry(sess)
        entry = next((k for k in registry['knowledgeBases'] if k.get('id') == kb_id), None)
        if not entry:
            return _kb_error('Knowledge base not found.', 404)
        existing = set(entry.get('fileList') or [])
        new_files = [p for p in stored if p not in existing]
        entry.setdefault('fileList', []).extend(new_files)
        if new_files:
            entry['indexState'] = {'status': 'pending', 'fileCount': 0, 'chunkCount': 0, 'lastIndexed': None}
        _save_kb_registry(sess, registry)
    if new_files:
        # Publish state before launching the worker so the frontend's first
        # refresh reliably starts its progress polling.
        _kb_progress[kb_id] = {'status': 'pending', 'done': 0, 'total': 0, 'error': None}
        threading.Thread(target=_kb_run_index, args=(sess['accountId'], kb_id, sess['key'], False), daemon=True).start()
    return _kb_ok({'added': len(new_files), 'rejected': rejected})


@app.before_request
def check_origin():
    if request.path.startswith('/proxy/') or request.path.startswith('/store') or request.path.startswith('/agent'):
        origin = request.headers.get('Origin', '')
        host   = request.headers.get('Host', '')
        if origin and origin not in ALLOWED_ORIGINS:
            return Response('{"error":"Origin not allowed."}',
                            403, content_type='application/json')
        if host and not (host.startswith('localhost:') or host.startswith('127.0.0.1:')):
            return Response('{"error":"Host not allowed."}',
                            403, content_type='application/json')

# ── Private IP ranges (SSRF protection) ──────────────────────────
# Three tiers:
#  - LOOPBACK_NETWORKS: same machine only. Always allowed (this is the mode
#    that already worked - local LM Studio/Ollama/vLLM on 'localhost').
#  - ALWAYS_BLOCKED_NETWORKS: never reachable via the proxy. Cloud metadata
#    endpoints and reserved/documentation/broadcast ranges - no legitimate
#    "my own LAN device" use case, only an SSRF one.
#  - LAN_NETWORKS: private/local-network ranges (RFC1918, link-local, IPv6
#    ULA, CGNAT). Blocked by default, unlockable per-provider via the
#    "kic_lan_confirm" marker after the user double-confirms the address
#    in the Provider editor - what makes an LM Studio/Ollama instance on
#    another PC on the network reachable.
LOOPBACK_NETWORKS = [
    ipaddress.ip_network('127.0.0.0/8'), ipaddress.ip_network('::1/128'),
]
ALWAYS_BLOCKED_NETWORKS = [
    ipaddress.ip_network('169.254.169.254/32'),  # cloud metadata (AWS/GCP/Azure/...)
    ipaddress.ip_network('0.0.0.0/8'),     ipaddress.ip_network('192.0.0.0/24'),
    ipaddress.ip_network('198.18.0.0/15'), ipaddress.ip_network('198.51.100.0/24'),
    ipaddress.ip_network('203.0.113.0/24'),ipaddress.ip_network('240.0.0.0/4'),
    ipaddress.ip_network('255.255.255.255/32'),
    ipaddress.ip_network('::/128'),        ipaddress.ip_network('2002::/16'),
    ipaddress.ip_network('100::/64'),      ipaddress.ip_network('64:ff9b::/96'),
]
LAN_NETWORKS = [
    ipaddress.ip_network('10.0.0.0/8'),    ipaddress.ip_network('172.16.0.0/12'),
    ipaddress.ip_network('192.168.0.0/16'),ipaddress.ip_network('169.254.0.0/16'),
    ipaddress.ip_network('100.64.0.0/10'), ipaddress.ip_network('fc00::/7'),
    ipaddress.ip_network('fe80::/10'),
]

def _resolve_all_ips(hostname):
    try:
        socket.setdefaulttimeout(5)
        return [i[4][0] for i in socket.getaddrinfo(hostname, None)]
    except Exception: return []
    finally: socket.setdefaulttimeout(None)

def _classify_ip(addr):
    if isinstance(addr, ipaddress.IPv6Address) and addr.ipv4_mapped:
        addr = addr.ipv4_mapped
    if any(addr in net for net in LOOPBACK_NETWORKS): return 'loopback'
    if any(addr in net for net in ALWAYS_BLOCKED_NETWORKS): return 'blocked'
    if any(addr in net for net in LAN_NETWORKS): return 'lan'
    return 'public'

def classify_host(host):
    """Classify a hostname or IP as 'loopback' | 'blocked' | 'lan' | 'public'
    | 'unresolvable', returning (class, ip) - the specific resolved IP the
    decision is based on. For hostnames, every resolved address is checked
    and the most restrictive class wins."""
    try:
        ipaddress.ip_address(host)
        ips = [host]
    except ValueError:
        ips = _resolve_all_ips(host)
    if not ips: return 'unresolvable', None
    classified = []
    for ip_str in ips:
        try: classified.append((ip_str, _classify_ip(ipaddress.ip_address(ip_str))))
        except ValueError: return 'blocked', None
    classes = {c for _, c in classified}
    if 'blocked' in classes: return 'blocked', None
    if 'lan' in classes:
        pin = next(ip for ip, c in classified if c == 'lan')
        return 'lan', pin
    if classes == {'loopback'}:
        return 'loopback', classified[0][0]
    if 'loopback' in classes:  # mixed loopback+public - treat as needing confirmation, be safe
        pin = next(ip for ip, c in classified if c == 'loopback')
        return 'lan', pin
    pin = next(ip for ip, c in classified if c == 'public')
    return 'public', pin

def is_allowed(target_url, method='GET', confirmed=False):
    # No domain allowlist here - users can point POST at any custom
    # OpenAI-compatible endpoint. SSRF protection: http/https only, and a
    # tiered check on where the address actually points (see above).
    # Returns (ok, reason, needs_confirm, pinned_ip). pinned_ip is the exact
    # address this decision was based on - the caller MUST connect to that
    # same address (see _pin_dns below) rather than letting requests/urllib3
    # re-resolve the hostname a second time, or a DNS answer that changes
    # between this check and the real connect (DNS rebinding) could steer
    # the actual request at a target this check never saw/approved.
    try: parsed = urlparse(target_url)
    except Exception: return False, 'Invalid URL', False, None
    if parsed.scheme not in ('http', 'https'): return False, 'HTTP/HTTPS only', False, None
    host = parsed.hostname or ''
    if not host: return False, 'No hostname', False, None
    # Fast-path for the exact strings the app has always used for "this
    # machine" - keeps behaving exactly as before even if DNS is weird.
    # No pinning needed: these always resolve locally, nothing to rebind to.
    if host in ('localhost', '127.0.0.1', '::1'):
        return True, '', False, None

    cls, pin = classify_host(host)
    if cls == 'unresolvable':
        return False, 'Host could not be resolved', False, None
    if cls == 'blocked':
        return False, 'This address is not allowed (reserved/blocked range).', False, None
    if cls == 'loopback':
        return True, '', False, pin
    if cls == 'lan':
        if not confirmed:
            return False, 'Local-network address - confirm it in the API panel first.', True, None
        return True, '', False, pin
    return True, '', False, pin  # public

# ── DNS pinning (closes the SSRF check's DNS-rebinding gap) ──────
# is_allowed()/classify_host() resolve and vet the hostname's IP, but the
# actual connection resolves it again - a low-TTL DNS answer could then
# point somewhere the check never saw. This pins socket.getaddrinfo() to
# the already-vetted address for that one hostname, for the request's
# duration.
#
# Uses thread-local state rather than a global save/restore swap: waitress
# serves multiple worker threads, and two concurrent "swap in, ..., swap
# back" calls on the same global `socket.getaddrinfo` would race. A single
# patched function that only special-cases the current thread's pin avoids
# that - unrelated hostnames/threads fall through to the real resolver.
_dns_pin = threading.local()
_real_getaddrinfo = socket.getaddrinfo

def _pinned_getaddrinfo(host, port, *args, **kwargs):
    if host == getattr(_dns_pin, 'host', None):
        ip = _dns_pin.ip
        family = socket.AF_INET6 if ':' in ip else socket.AF_INET
        return [(family, socket.SOCK_STREAM, socket.IPPROTO_TCP, '', (ip, port))]
    return _real_getaddrinfo(host, port, *args, **kwargs)

socket.getaddrinfo = _pinned_getaddrinfo

@contextlib.contextmanager
def _pin_dns(host, ip):
    """Pin `host` to `ip` for the current thread for the duration of the
    'with' block. No-op if ip is None (nothing to pin - e.g. the
    localhost/127.0.0.1/::1 fast path, which needs no rebinding protection)."""
    if not ip:
        yield
        return
    prev_host, prev_ip = getattr(_dns_pin, 'host', None), getattr(_dns_pin, 'ip', None)
    _dns_pin.host, _dns_pin.ip = host, ip
    try:
        yield
    finally:
        _dns_pin.host, _dns_pin.ip = prev_host, prev_ip

# ── Rate-Limiting ─────────────────────────────────────────────────
_rate_data = defaultdict(lambda: {'count': 0, 'reset': 0.0})
_rate_lock = threading.Lock()
RATE_LIMIT = 120; RATE_WINDOW = 60.0
_CLEANUP_EVERY = 300; _last_cleanup = time.monotonic()

def check_rate_limit(ip):
    global _last_cleanup
    now = time.monotonic()
    with _rate_lock:
        if now - _last_cleanup > _CLEANUP_EVERY:
            for k in [k for k, v in _rate_data.items() if now > v['reset']]:
                del _rate_data[k]
            _last_cleanup = now
        d = _rate_data[ip]
        if now > d['reset']: d['count'] = 0; d['reset'] = now + RATE_WINDOW
        d['count'] += 1
        return d['count'] <= RATE_LIMIT

# ── CORS + Security Headers ───────────────────────────────────────
CORS_HEADERS = {
    'Access-Control-Allow-Origin':  f'http://localhost:{port}',
    'Access-Control-Allow-Methods': 'GET, POST, PUT, DELETE, OPTIONS',
    'Access-Control-Allow-Headers': (
        'Authorization, Content-Type, x-api-key, '
        'X-Subscription-Token, xi-api-key, '
        'anthropic-version, anthropic-dangerous-direct-browser-access, '
        'HTTP-Referer, X-Title, '
        'Ocp-Apim-Subscription-Key'           # Bing Search API
    ),
}
EXCLUDED_RESP_HEADERS = {
    # Hop-by-hop headers (RFC 7230 §6.1) - forwarding these violates PEP 3333
    # and crashes waitress.
    'connection','keep-alive','proxy-authenticate','proxy-authorization',
    'te','trailer','trailers','transfer-encoding','upgrade',
    'content-encoding','content-length','server','x-powered-by','set-cookie','location',
}
SECURITY_HEADERS = {
    'X-Content-Type-Options': 'nosniff',
    'X-Frame-Options': 'DENY',
    'Referrer-Policy': 'no-referrer',
    'X-DNS-Prefetch-Control': 'off',
    'Strict-Transport-Security': 'max-age=31536000; includeSubDomains',
    'Content-Security-Policy': (
        # Actual enforced policy - the <meta> CSP in kiconnect.html is kept
        # in sync but this header wins if the two ever disagree. No
        # 'unsafe-inline' needed since former inline scripts now live in
        # kiconnect-mathjax-config.js / kiconnect.js.
        "default-src 'self'; "
        "script-src 'self'; "
        "style-src 'self' 'unsafe-inline'; "
        "connect-src 'self' https://api.anthropic.com https://api.openai.com "
        "https://chat.kiconnect.nrw https://openrouter.ai "
        "https://api.mistral.ai https://generativelanguage.googleapis.com "
        "https://texttospeech.googleapis.com "
        "https://api.x.ai https://api.groq.com "
        "https://api.deepseek.com https://api.minimax.io "
        "https://api.z.ai https://api.moonshot.ai "
        "https://api.elevenlabs.io "
        "https://api.search.brave.com https://html.duckduckgo.com "
        "https://lite.duckduckgo.com https://api.qwant.com https://search.yahoo.com "
        "https://www.startpage.com https://www.googleapis.com https://api.bing.microsoft.com "
        "https://api.mojeek.com https://yandex.com "
        "https://api.langsearch.com "
        "https://searx.be https://searxng.world https://search.bus-hit.me "
        "https://searx.tiekoetter.com https://search.sapti.me https://searx.prvcy.eu "
        "https://searx.fmac.xyz https://search.ononoki.org; "
        "img-src 'self' data: blob:; "
        # TTS providers (OpenAI/ElevenLabs/Groq) return audio bytes that get
        # played back via `new Audio(URL.createObjectURL(blob))` — needs
        # media-src to allow blob:, otherwise it silently falls back to
        # default-src 'self' and playback is blocked.
        "media-src 'self' blob:; "
        "font-src 'self'; "
        "worker-src 'self' blob:; "
        "frame-src 'none'; object-src 'none'; base-uri 'self';"
    ),
    'Permissions-Policy': (
        'geolocation=(), camera=(), '
        'payment=(), usb=(), magnetometer=(), gyroscope=()'
    ),
}

@app.after_request
def add_security_headers(response):
    for k, v in CORS_HEADERS.items():   response.headers[k] = v
    for k, v in SECURITY_HEADERS.items(): response.headers[k] = v
    return response

# ── Statische Dateien ─────────────────────────────────────────────
@app.route('/')
def index():
    return send_from_directory(STATIC_DIR, 'kiconnect.html')

@app.route('/<path:filename>')
def static_files(filename):
    if filename.startswith('proxy/'):
        return _proxy_request(filename[len('proxy/'):])
    safe_path = os.path.realpath(os.path.join(STATIC_DIR, filename))
    if not safe_path.startswith(STATIC_DIR + os.sep) and safe_path != STATIC_DIR:
        abort(404)
    basename = os.path.basename(filename).lower()
    if basename.startswith('.') or any(
        basename.endswith(e) for e in ['.py','.key','.pem','.env','.ini','.cfg','.log']
    ): abort(404)
    if os.path.isfile(safe_path):
        return send_from_directory(STATIC_DIR, filename)
    abort(404)

# ── Proxy ─────────────────────────────────────────────────────────
@app.route('/proxy/<path:target>', methods=['GET', 'POST', 'OPTIONS'])
def proxy(target):
    if request.method == 'OPTIONS':
        return Response('', 204, headers={**CORS_HEADERS, **SECURITY_HEADERS})
    return _proxy_request(target)

def _proxy_request(target_url):
    client_ip = request.remote_addr or '0.0.0.0'
    if not check_rate_limit(client_ip):
        return Response('{"error":"Too many requests - please wait a moment."}',
                        429, content_type='application/json')
    try: target_url = unquote(target_url)
    except Exception: pass

    # "kic_lan_confirm=1" is a marker the frontend appends to the proxy URL
    # (never to the upstream one) only for a provider address the user has
    # explicitly double-confirmed in the Provider editor - see
    # confirmLanAddress() in kiconnect.js. Strip it before forwarding so it
    # never reaches the upstream API as a stray query parameter.
    fwd_params = request.args.copy()
    lan_confirmed = fwd_params.pop('kic_lan_confirm', None) == '1'

    ok, reason, needs_confirm, pinned_ip = is_allowed(target_url, request.method, confirmed=lan_confirmed)
    if not ok:
        print(f'  blocked [{reason}]')
        status = 428 if needs_confirm else 403
        body = json.dumps({'error': reason, 'code': 'lan_confirm_required' if needs_confirm else 'blocked'})
        return Response(body, status, content_type='application/json')

    ALLOWED_REQ_HEADERS = {
        'authorization','content-type','x-api-key',
        'x-subscription-token',
        'xi-api-key',                          # ElevenLabs TTS auth
        'anthropic-version','anthropic-dangerous-direct-browser-access','anthropic-beta',
        'accept','http-referer','x-title','origin',
        'ocp-apim-subscription-key',          # Bing Search API
        'user-agent',                          # Search engines (browser can't send it, proxy injects below)
        'referer','accept-language','sec-fetch-site','sec-fetch-mode','sec-fetch-dest',
    }
    fwd_headers = {k: v for k, v in request.headers if k.lower() in ALLOWED_REQ_HEADERS}
    # Inject a browser-like User-Agent when none is present - required by DuckDuckGo,
    # Bing, Brave etc. (browsers cannot set User-Agent in fetch, so the proxy must do it)
    if not any(k.lower() == 'user-agent' for k in fwd_headers):
        fwd_headers['User-Agent'] = (
            'Mozilla/5.0 (Windows NT 10.0; Win64; x64) '
            'AppleWebKit/537.36 (KHTML, like Gecko) '
            'Chrome/124.0.0.0 Safari/537.36'
        )
    fwd_headers.setdefault('Accept-Language', 'de-DE,de;q=0.9,en-US;q=0.8,en;q=0.7')
    fwd_headers.setdefault('Accept', 'text/html,application/xhtml+xml,application/xml;q=0.9,application/json;q=0.8,*/*;q=0.7')
    if 'http-referer' in {k.lower() for k in fwd_headers} and not any(k.lower() == 'referer' for k in fwd_headers):
        for k, v in list(fwd_headers.items()):
            if k.lower() == 'http-referer':
                fwd_headers['Referer'] = v
                break
    body = request.get_data()
    if len(body) > MAX_BODY_SIZE:
        return Response('{"error":"Request body too large."}', 413, content_type='application/json')

    print(f'  -> {request.method:6s} {target_url[:90]}')
    target_host = urlparse(target_url).hostname or ''
    try:
        # Pinned to the exact IP is_allowed() already vetted above - see
        # _pin_dns()/_pinned_getaddrinfo() for why this matters.
        with _pin_dns(target_host, pinned_ip):
            upstream = requests.request(
                method=request.method, url=target_url, headers=fwd_headers,
                data=body, params=fwd_params, stream=True,
                timeout=(10, 180), verify=True, allow_redirects=False,
            )
        resp_headers = {k: v for k, v in upstream.headers.items()
                        if k.lower() not in EXCLUDED_RESP_HEADERS}
        resp_headers.update(CORS_HEADERS); resp_headers.update(SECURITY_HEADERS)

        def generate():
            try:
                for chunk in upstream.iter_content(chunk_size=8192):
                    if chunk: yield chunk
            except Exception as e:
                print(f'  Stream error: {type(e).__name__}')

        print(f'  <- {upstream.status_code}')
        return Response(generate(), status=upstream.status_code, headers=resp_headers)

    except requests.exceptions.SSLError:
        return Response('{"error":"SSL certificate invalid."}', 502, content_type='application/json')
    except requests.exceptions.ConnectionError:
        return Response('{"error":"Target server unreachable."}', 502, content_type='application/json')
    except requests.exceptions.Timeout:
        return Response('{"error":"API did not respond (timeout)."}', 504, content_type='application/json')
    except Exception as e:
        print(f'  Error: {type(e).__name__}')
        return Response('{"error":"Internal proxy error."}', 500, content_type='application/json')


if __name__ == '__main__':
    try:
        from waitress import serve
    except ImportError:
        print('[ERROR] waitress not installed. Run: pip install waitress')
        sys.exit(1)

    W = 72   # Frame len
    IW = W - 4  # Inner frame content width (68) — muss zur Randbreite von 6 passen

    def line(text=''):
        print('║ ' + text.ljust(W) + ' ║')

    def iline(text=''):
        print('║   ' + text.ljust(IW) + '   ║')

    print('╔' + '═' * (W + 2) + '╗')
    line('KI Connect - CORS-Proxy + Storage-Server (Waitress)')
    print('╠' + '═' * (W + 2) + '╣')
    line(f'Running on: http://localhost:{port}   Data dir: ./datas/')
    line()
    line('Storage-API (localhost only):')
    line('  GET/PUT        /store/           Account registry')
    line('  GET            /store/<id>       Keys list')
    line('  GET/PUT/DELETE /store/<id>/<k>   Data read/write')
    line()
    line('Agent-API (per-account encrypted project registry,')
    line('           requires an unlocked session):')
    print('║  ┌' + '─' * IW + '┐  ║')
    iline('POST     /agent/session/unlock|lock|rekey Unlock/drop/rekey session')
    iline('GET      /agent/browse                   Browse OS folders (+files=1 to list files)')
    iline('GET/POST /agent/projects                 List/register project')
    iline('DELETE   /agent/projects/<id>            Unregister (keeps files)')
    iline('PUT      /agent/projects/<id>/shell|path Toggle shell / change path')
    iline('POST     /agent/exec/<id>                Run shell command')
    iline('GET      /agent/tree|search/<id>         File listing / text search')
    iline('GET/PUT/DELETE /agent/file/<id>/<p>      Read/write/delete file')
    iline('POST/DELETE /agent/dir/<id>/<p>          Create/delete folder')
    iline('POST     /agent/move/<id>                Move/rename')
    iline('POST     /agent/copy/<id>                Copy (leaves original)')
    print('║  └' + '─' * IW + '┘  ║')
    line()
    line('Knowledge-Base API (RAG - reuses the Agent session above):')
    print('║  ┌' + '─' * IW + '┐  ║')
    iline('POST     /kb/create                      Register folder or file list, start indexing')
    iline('GET      /kb/list                        List knowledge bases')
    iline('GET      /kb/<id>/status                 Indexing progress')
    iline('POST     /kb/<id>/reindex                Re-index (full, v1)')
    iline('POST     /kb/<id>/add-files              Add individual files (incremental)')
    iline('DELETE   /kb/<id>                        Unregister (keeps files)')
    iline('POST     /kb/<id>/search                 Top-k similarity search')
    iline('PATCH    /kb/<id>/settings               Embedding/chunking config')
    print('║  └' + '─' * IW + '┘  ║')
    line()
    line('Proxy allowlist: anthropic, openai, openrouter, mistral, googleapis,')
    line('                 x.ai, groq, deepseek, minimax, z.ai, moonshot,')
    line('                 chat.kiconnect.nrw')
    line('Search: brave, duckduckgo, qwant, yahoo, startpage, google, bing,')
    line('        mojeek, yandex, searxng, langsearch')
    line()
    line('Stop: Ctrl+C')
    print('╚' + '═' * (W + 2) + '╝')

    threading.Timer(1.2, lambda: webbrowser.open(f'http://localhost:{port}')).start()
    serve(app, host='127.0.0.1', port=port, threads=8, channel_timeout=120, cleanup_interval=10)
